import logging
import re
from pptx import Presentation
from pptx.dml.color import RGBColor
from jinja2 import Environment, BaseLoader
import pandas as pd
from pptx.dml.color import MSO_COLOR_TYPE

# ---------- Emoji removal helper ----------
# Regex to match emoji and a broad set of pictographs/symbols.
_EMOJI_REGEX = re.compile(
    "["
    "\U0001f600-\U0001f64f"  # emoticons
    "\U0001f300-\U0001f5ff"  # symbols & pictographs
    "\U0001f680-\U0001f6ff"  # transport & map symbols
    "\U0001f1e0-\U0001f1ff"  # flags (iOS)
    "\U00002702-\U000027b0"  # dingbats
    "\U000024c2-\U0001f251"
    "\U0001f900-\U0001f9ff"  # Supplemental Symbols and Pictographs
    "\U0001fa70-\U0001faff"  # Symbols and Pictographs Extended-A
    "\U00002600-\U000026ff"  # Misc symbols
    "]+",
    flags=re.UNICODE,
)


def remove_emojis_from_string(s):
    if not isinstance(s, str):
        return s
    return _EMOJI_REGEX.sub("", s)


_WS_RE = re.compile(r"\s+")


def sanitize_header_name(h):
    """
    Remove emojis, strip, and replace inner whitespace with underscores.
    """
    if h is None:
        return h
    # convert to str in case it's not already
    s = str(h)
    # remove emojis and trailing/leading spaces handled by that function
    s = remove_emojis_from_string(s)
    # collapse internal whitespace to single underscore
    s = _WS_RE.sub("_", s)
    # strip leading/trailing underscores produced by replacement of leading/trailing spaces
    s = s.strip("_")
    return s


def sanitize_dataframe_headers(df):
    """
    Return a new DataFrame with sanitized column headers:
    - emojis removed
    - internal whitespace replaced with underscores
    - leading/trailing underscores trimmed
    """
    # build rename mapping
    rename_map = {col: sanitize_header_name(col) for col in df.columns}
    # return a copy with renamed columns
    return df.rename(columns=rename_map)


# ---------- Marker encoding for colored segments ----------
_START = "\u0002"
_SEP = "\u0003"
_END = "\u0004"
MARKER_RE = re.compile(
    re.escape(_START)
    + r"(#[0-9A-Fa-f]{6})"
    + re.escape(_SEP)
    + r"(.*?)"
    + re.escape(_END),
    re.DOTALL,
)


def encode_colored_text(text, hex_color):
    return f"{_START}{hex_color}{_SEP}{text}{_END}"


def decode_marked_segments(rendered_text):
    segments = []
    pos = 0
    for m in MARKER_RE.finditer(rendered_text):
        if m.start() > pos:
            segments.append((rendered_text[pos : m.start()], None))
        hex_color = m.group(1)
        txt = m.group(2)
        segments.append((txt, hex_color))
        pos = m.end()
    if pos < len(rendered_text):
        segments.append((rendered_text[pos:], None))
    return segments


# ---------- Formatting copy helper ----------
def copy_run_format(src_run, dest_run):
    try:
        dest_run.font.bold = src_run.font.bold
        dest_run.font.italic = src_run.font.italic
        dest_run.font.underline = src_run.font.underline
        dest_run.font.strike = src_run.font.strike
        if src_run.font.name:
            dest_run.font.name = src_run.font.name
        if src_run.font.size:
            dest_run.font.size = src_run.font.size
        # Try to copy RGB color if available
        try:
            col = src_run.font.color
            if (
                getattr(col, "type", None) is not None
                and getattr(col, "rgb", None) is not None
            ):
                rgb = col.rgb
                try:
                    dest_run.font.color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
                except Exception:
                    try:
                        dest_run.font.color.rgb = rgb
                    except Exception:
                        pass
        except Exception:
            pass
    except Exception:
        pass


# ---------- Robust colorize_positive filter ----------
def colorize_positive(
    value, positive_hex="#008000", negative_hex="#C00000", zero_hex="#000000"
):
    """
    Try to parse 'value' robustly; return marker-wrapped text for coloring.
    """

    def try_parse_number(v):
        if v is None:
            return None
        if isinstance(v, (int, float)):
            try:
                return float(v)
            except Exception:
                return None
        # pandas NA / numpy.nan
        try:
            import numpy as _np

            if v is _np.nan:
                return None
        except Exception:
            pass
        if isinstance(v, str):
            s = v.strip()
            if s == "":
                return None
            # Remove emojis before parsing (safety)
            s = remove_emojis_from_string(s)
            s = s.replace(",", "").replace(" ", "")
            s = re.sub(r"^[^\d\-\+\.]+", "", s)
            s = re.sub(r"[^\d\.eE\-\+]$", "", s)
            try:
                return float(s)
            except Exception:
                m = re.match(r"^\(([\d\.,\-]+)\)$", v.strip())
                if m:
                    inner = m.group(1).replace(",", "")
                    try:
                        return -float(inner)
                    except Exception:
                        return None
                return None
        try:
            return float(v)
        except Exception:
            return None

    num = try_parse_number(value)
    if num is None:
        hex_color = zero_hex
    elif num > 0:
        hex_color = positive_hex
    elif num < 0:
        hex_color = negative_hex
    else:
        hex_color = zero_hex

    text = "" if value is None else str(value)
    # ensure the output text also has emojis removed (defensive)
    text = remove_emojis_from_string(text)
    return encode_colored_text(text, hex_color)


# ---------- Jinja env ----------
def make_jinja_env():
    env = Environment(loader=BaseLoader(), autoescape=False)
    env.filters["colorize_positive"] = colorize_positive
    return env


def render_text_with_jinja(text, context, env=None):
    if env is None:
        env = make_jinja_env()
    template = env.from_string(text)
    return template.render(**(context or {}))


# ---------- Extract original runs and text ----------
def extract_text_and_run_meta(text_frame):
    parts = []
    run_meta = []
    for p in text_frame.paragraphs:
        for r in p.runs:
            run_text = r.text or ""
            parts.append(run_text)
            run_meta.append({"text": run_text, "run_obj": r})
        parts.append("\n")
        run_meta.append({"text": "\n", "run_obj": None})
    if parts and parts[-1] == "\n":
        parts.pop()
        run_meta.pop()
    full_text = "".join(parts)
    return full_text, run_meta


# ---------- High-level processor ----------
def process_text_field(shape, text_to_insert, df, env=None):
    text_to_insert = str(text_to_insert)
    jinja_tag_re = re.compile(r"({{.*?}}|{%.+?%})", re.DOTALL)
    text_frame = shape.text_frame
    full_text, run_meta = extract_text_and_run_meta(text_frame)

    if not jinja_tag_re.search(full_text):
        logging.debug("No Jinja tags found in shape; applying fallback if different.")
        if full_text != (text_to_insert or ""):
            update_text_frame_preserving_formatting(text_frame, text_to_insert or "")
        return

    df_sanitized = sanitize_dataframe_headers(df)
    rows = df_sanitized.to_dict(orient="records")
    context = {"rows": rows}

    rendered = render_text_with_jinja(full_text, context, env=env)

    # --- Compare old vs new to decide whether to modify ---
    if rendered.strip() == full_text.strip():
        logging.debug("Rendered Jinja output identical to original; skipping update.")
        return

    # --- Update text safely ---
    reinsert_rendered_text_preserving_formatting(text_frame, rendered, run_meta)


def update_text_frame_preserving_formatting(text_frame, new_text):
    """
    Replace text content but preserve shape formatting and paragraph style.
    """
    # Grab formatting from the first run
    if not text_frame.paragraphs:
        text_frame.text = new_text
        return

    p = text_frame.paragraphs[0]
    runs = p.runs
    font = runs[0].font if runs else None
    if font and getattr(font, "color", None):
        col = font.color
        try:
            if getattr(col, "type", None) == MSO_COLOR_TYPE.RGB and getattr(
                col, "rgb", None
            ):
                color = col.rgb
            elif getattr(col, "type", None) == MSO_COLOR_TYPE.SCHEME and getattr(
                col, "theme_color", None
            ):
                color = col.theme_color
        except Exception:
            pass
    size = font.size if font and font.size else Pt(12)

    # Clear all text but keep paragraphs
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.text = ""

    # Replace only first run text (preserves style)
    if not text_frame.paragraphs:
        p = text_frame.add_paragraph()

    if not p.runs:
        run = p.add_run()
    else:
        run = p.runs[0]

    run.text = new_text

    # Reapply original font attributes (if any)
    if color:
        try:
            if isinstance(color, RGBColor):
                run.font.color.rgb = color
            else:
                run.font.color.theme_color = color
        except Exception:
            pass
    if size:
        run.font.size = size


def reinsert_rendered_text_preserving_formatting(
    text_frame, rendered_text, run_meta=None
):
    first_paragraph = text_frame.paragraphs[0] if text_frame.paragraphs else None
    first_run = (
        first_paragraph.runs[0] if first_paragraph and first_paragraph.runs else None
    )
    font = getattr(first_run, "font", None)
    alignment = getattr(first_paragraph, "alignment", None)

    for p in list(text_frame.paragraphs):
        text_frame._element.remove(p._p)

    new_paragraph = text_frame.add_paragraph()
    new_run = new_paragraph.add_run()
    new_run.text = rendered_text

    # Safely copy all style attributes
    copy_font_format(font, new_run.font)

    if alignment is not None:
        new_paragraph.alignment = alignment


def copy_font_format(src_font, dest_font):
    """Copy color, size, bold, italic safely between fonts."""
    if not src_font or not dest_font:
        return

    color = src_font.color
    if color and color.type is not None:
        if color.type == MSO_COLOR_TYPE.RGB and color.rgb:
            dest_font.color.rgb = color.rgb
        elif color.type == MSO_COLOR_TYPE.SCHEME and color.theme_color is not None:
            dest_font.color.theme_color = color.theme_color

    if src_font.size:
        dest_font.size = src_font.size
    if src_font.bold is not None:
        dest_font.bold = src_font.bold
    if src_font.italic is not None:
        dest_font.italic = src_font.italic
