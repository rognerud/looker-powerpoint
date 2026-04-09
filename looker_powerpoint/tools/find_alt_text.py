from pptx import Presentation
from lxml import etree
import yaml

NS = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}

# Mapping of "smart" / non-ASCII quote variants to their ASCII equivalents.
_QUOTE_REPLACEMENTS = {
    # Double-quote variants → straight double quote
    "\u201c": '"',  # LEFT DOUBLE QUOTATION MARK  "
    "\u201d": '"',  # RIGHT DOUBLE QUOTATION MARK "
    "\u201e": '"',  # DOUBLE LOW-9 QUOTATION MARK „
    "\u2033": '"',  # DOUBLE PRIME               ″
    "\u00ab": '"',  # LEFT-POINTING DOUBLE ANGLE  «
    "\u00bb": '"',  # RIGHT-POINTING DOUBLE ANGLE »
    # Single-quote / apostrophe variants → straight single quote
    "\u2018": "'",  # LEFT SINGLE QUOTATION MARK  '
    "\u2019": "'",  # RIGHT SINGLE QUOTATION MARK '
    "\u201a": "'",  # SINGLE LOW-9 QUOTATION MARK ‚
    "\u2032": "'",  # PRIME                       ′
    "\u0060": "'",  # GRAVE ACCENT                `
    "\u00b4": "'",  # ACUTE ACCENT                ´
}

# Build a translation table once for efficient single-pass replacement.
_QUOTE_TABLE = str.maketrans(_QUOTE_REPLACEMENTS)


def cleanse_alt_text(text: str) -> str:
    """Normalise alternative-text before YAML parsing.

    Replaces typographic / "smart" quote characters with their plain ASCII
    equivalents so that YAML produced by applications that substitute curly
    quotes (e.g. macOS, Word, PowerPoint) can still be parsed correctly.

    Args:
        text: Raw alternative-text string extracted from a shape.

    Returns:
        The cleansed string with all known fancy quote variants replaced by
        straight ASCII quotes.
    """
    return text.translate(_QUOTE_TABLE)


def extract_alt_text(shape):
    """
    Extracts the alternative text description from a shape's XML.

    Args:
        shape: A Shape object from pptx.

    Returns:
        str: The alternative text description, or None if not found.
    """
    xml_str = shape.element.xml  # get XML string of the shape element
    xml_elem = etree.fromstring(xml_str)  # parse it into an lxml element
    for path in [
        ".//p:nvSpPr/p:cNvPr",
        ".//p:nvPicPr/p:cNvPr",
        ".//p:nvGraphicFramePr/p:cNvPr",
    ]:
        cNvPr_elements = xml_elem.xpath(path, namespaces=NS)
        if cNvPr_elements:
            descr = cNvPr_elements[0].get("descr")
            if descr:
                data = yaml.safe_load(
                    cleanse_alt_text(descr)
                )  # Use safe_load for untrusted sources

                return data
    return None


def get_presentation_objects_with_descriptions(pptx_path):
    """
    Extracts all shapes from a PowerPoint presentation and returns them with descriptions.

    Args:
        pptx_path (str): The path to the PowerPoint presentation file.

    Returns:
        list: A list of dictionaries, where each dictionary represents a shape and
              contains the shape object, its description, and the slide number.
              Returns an empty list if the presentation cannot be opened or has no slides/shapes.
    """
    try:
        presentation = Presentation(pptx_path)
    except Exception as e:
        print(f"Error opening presentation: {e}")
        return []

    objects_with_descriptions = []

    for i, slide in enumerate(presentation.slides, start=0):
        for shape in slide.shapes:
            description = extract_alt_text(shape)  # Generate description

            emu_to_pixels = lambda emu: emu / 9525

            width_px = emu_to_pixels(shape.width)
            height_px = emu_to_pixels(shape.height)

            if description:
                if type(description) is dict and "meta_name" in description:
                    shape_id = description.get("meta_name")
                else:
                    shape_id = (
                        f"{i},{shape.shape_id}"  # Use shape number for identification
                    )

                objects_with_descriptions.append(
                    {
                        "shape_id": shape_id,  # Use shape number for identification
                        "shape_type": shape.shape_type.name,
                        "shape_width": round(width_px),
                        "shape_height": round(height_px),
                        "integration": description,
                        "slide_number": i,  # Use the enumerate index for slide number
                        "shape_number": shape.shape_id,
                    }
                )

    return objects_with_descriptions


if __name__ == "__main__":
    # Example Usage
    pptx_file = "p.pptx"  # Replace with your file path
    objects = get_presentation_objects_with_descriptions(pptx_file)
    from rich import print
