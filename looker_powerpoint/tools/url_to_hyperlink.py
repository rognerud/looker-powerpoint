import re

from pptx.dml.color import RGBColor


def add_text_with_numbered_links(text_frame, text, start_index=1):
    """
    Replaces URLs in `text` with numbered references "(1)", "(23)", etc.
    The placeholder is hyperlinked to the URL.
    - Clears any prior content in the text_frame.
    - Removes newlines if hyperlinks are present.
    - If a URL ends with digits, uses that number instead of auto numbering.
    Returns the next available numeric index.
    """
    text_frame.clear()

    url_pattern = re.compile(r"https?://\S+")
    paragraph = text_frame.add_paragraph()
    index = start_index

    matches = url_pattern.findall(text)
    if matches:
        # Flatten newlines if hyperlinks exist
        text = text.replace("\n", " ")

    # Split while keeping URLs in the list
    parts = re.split(f"({url_pattern.pattern})", text)

    for part in parts:
        if not part:
            continue

        if url_pattern.fullmatch(part.strip()):
            url = part.strip()

            # Detect digits at the end of the URL path
            match_digits = re.search(r"(\d+)(?:[/?#]?)*$", url)
            number_text = match_digits.group(1) if match_digits else str(index)

            run = paragraph.add_run()
            run.text = f"({number_text})"
            run.hyperlink.address = url
            run.font.color.rgb = RGBColor(0, 0, 255)
            run.font.underline = True

            index += 1 if not match_digits else 0
        else:
            run = paragraph.add_run()
            run.text = part

    return index
