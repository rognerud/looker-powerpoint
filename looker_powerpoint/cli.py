import argparse
import asyncio
import collections
import datetime
import io
import json
import logging
import os
import re
import subprocess
from io import BytesIO

import pandas as pd
from lxml import etree
from PIL import Image
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pydantic import ValidationError
from rich.logging import RichHandler
from rich_argparse import RichHelpFormatter
import requests

from looker_powerpoint import gemini as gemini_module
from looker_powerpoint.looker import LookerClient
from looker_powerpoint.models import LookerShape, GeminiShape
from looker_powerpoint.tools.find_alt_text import (
    get_presentation_objects_with_descriptions,
)
from looker_powerpoint.tools.pptx_text_handler import (
    process_text_field,
    update_text_frame_preserving_formatting,
)

NS = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}


class Cli:
    # color with rich
    HEADER = """
        Looker PowerPoint CLI :
        A command line interface for Looker PowerPoint integration.
    """

    def __init__(self):
        self.client = None
        self.relevant_shapes = []
        self.looker_shapes = []
        self.gemini_shapes = []
        self.data = {}

        # Initialize the argument parser
        self.parser = self._init_argparser()

        # load tools
        self.get_alt_text = get_presentation_objects_with_descriptions

    def _init_looker(self):
        """Initialize the Looker client"""
        if not self.args.debug_queries:
            logging.getLogger("looker_sdk").setLevel(logging.ERROR)
        self.client = LookerClient()

    def _init_argparser(self):
        """Create and configure the argument parser"""
        parser = argparse.ArgumentParser(
            description=self.HEADER,
            formatter_class=RichHelpFormatter,
        )
        # todo
        parser.add_argument(
            "--file-path",
            "-f",
            help="Path to the PowerPoint file to process",
            default=None,
            type=str,
        )
        # todo
        parser.add_argument(
            "--output-dir",
            "-o",
            help="""Path to a directory that will contain the generated pptx files. \n
                .env: OUTPUT_DIR""",
            default="output",
            type=str,
        )
        # todo
        parser.add_argument(
            "--add-links",
            help="Add links to looker in the slides. \n .env: ADD_LINKS",
            action="store_true",
            default=False,
        )
        # todo
        parser.add_argument(
            "--hide-errors",
            help="""
                Stop showing red outlines around shapes with errors. \n
                .env: HIDE_ERRORS
            """,
            action="store_true",
            default=False,
        )
        # todo
        parser.add_argument(
            "--parse-date-syntax-in-filename",
            "-p",
            help="""Parse date syntax in the filename. \n
                .env: PARSE_DATE_SYNTAX_IN_FILENAME
                """,
            action="store_true",
            default=True,
        )

        parser.add_argument(
            "--self",
            "-s",
            help="""Replace the powerpoint file directly instead of creating a new file. \n
                .env: SELF""",
            action="store_true",
            default=False,
        )

        parser.add_argument(
            "--quiet",
            "-q",
            help="""Do not open the PowerPoint file after processing. \n
                .env: QUIET""",
            action="store_true",
            default=False,
        )

        parser.add_argument(
            "--filter",
            help="""use the string to filter shapes if they have a set filter dimension""",
            action="store",
            default=None,
            type=str,
        )

        parser.add_argument(
            "--debug-queries",
            help="""Enable debugging for Looker queries. \n
                .env: DEBUG_QUERIES""",
            action="store_true",
            default=False,
        )

        parser.add_argument(
            "-v",
            "--verbose",
            action="count",
            default=0,
            help="Increase verbosity (e.g., -v, -vv, -vvv)",
        )

        return parser

    def _setup_logging(self):
        if self.args.verbose == 0:
            level = logging.WARNING
        elif self.args.verbose == 1:
            level = logging.INFO
        else:
            level = logging.DEBUG

        logging.basicConfig(
            level=level,
            format="%(message)s",
            datefmt="[%X]",
            handlers=[RichHandler()],
        )

    def _pick_file(self):
        """
        Picks the PowerPoint file to process.
        If no file path is provided, it looks for the first .pptx file in the current directory.

        Returns:
            str: The path to the PowerPoint file.
        """
        self.file_path = self.args.file_path

        if self.file_path:
            try:
                self.presentation = Presentation(self.file_path)
            except Exception as e:
                logging.error(f"Error opening {self.file_path}: {e}")
        else:
            # If no file path is provided look for a file in the current directory
            files = [
                f
                for f in os.listdir(".")
                if f.endswith(".pptx") and not f.startswith("~$")
            ]
            if files:
                self.file_path = files[0]
                logging.warning(
                    f"No file path provided, using first found file: {self.file_path}. To specify a file, use the -f flag like 'lpt -f <file_path>'."
                )

                try:
                    self.presentation = Presentation(self.file_path)
                except Exception as e:
                    logging.error(f"Error opening {self.file_path}: {e}")
            else:
                logging.error(
                    """
                    No PowerPoint file found in the current directory, closing.
                    Run from a directory with a .pptx file, or
                    specify file using -f flag like 'lpt -f <file_path>'.
                """
                )
                exit(1)

    def _fill_table(self, table, df, headers=True):
        """
        Fills a PowerPoint table with data from a DataFrame.

        Args:
            table: A Table object from pptx.
            df: A pandas DataFrame containing the data to fill the table.
        """

        # Get table dimensions
        table_rows = len(table.rows)
        table_cols = len(table.columns)

        # Get DataFrame dimensions
        df_rows = df.shape[0] + 1  # +1 for header
        df_cols = df.shape[1]

        # Determine how much we can fill
        rows_to_fill = min(table_rows, df_rows)
        cols_to_fill = min(table_cols, df_cols)

        # Fill header row
        if headers:
            for col_idx in range(cols_to_fill):
                table.cell(0, col_idx).text = str(df.columns[col_idx])

        # Fill DataFrame values
        for row_idx in range(1, rows_to_fill):  # skip header row
            for col_idx in range(cols_to_fill):
                value = df.iloc[row_idx - 1, col_idx]
                table.cell(row_idx, col_idx).text = str(value)

        # Optional: Clear unused cells
        for row_idx in range(rows_to_fill, table_rows):
            for col_idx in range(table_cols):
                table.cell(row_idx, col_idx).text = ""

        for col_idx in range(cols_to_fill, table_cols):
            for row_idx in range(table_rows):
                table.cell(row_idx, col_idx).text = ""

    def _set_alt_text(self, shape, data):
        """
        Sets the alternative text description for a shape's XML.

        Args:
            shape: A Shape object from pptx.
            data: A Python object (dict, list, etc.) to serialize and set as YAML in the descr attribute.
        """
        xml_str = shape.element.xml
        xml_elem = etree.fromstring(xml_str)
        import yaml

        # convert pydantic model to dict
        if isinstance(data, dict) is False:
            data = data.model_dump()
        data = {k: v for k, v in data.items() if v is not None}
        data = yaml.dump(data)

        # remove None values from data, and convert to string with newlines for YAML compatibility

        for path in [
            ".//p:nvSpPr/p:cNvPr",
            ".//p:nvPicPr/p:cNvPr",
            ".//p:nvGraphicFramePr/p:cNvPr",
        ]:
            cNvPr_elements = xml_elem.xpath(path, namespaces=NS)
            if cNvPr_elements:
                cNvPr = cNvPr_elements[0]
                yaml_text = str(data)
                cNvPr.set("descr", yaml_text)

                # Overwrite the element in the actual pptx shape with updated XML
                shape_element = shape.element
                new_element = etree.fromstring(etree.tostring(xml_elem))
                shape_element.clear()
                for child in new_element:
                    shape_element.append(child)
                return

        raise ValueError("No compatible cNvPr element found to set descr.")

    def _mark_failure(self, slide, shape):
        line_color_rgb = (255, 0, 0)  # RGB color for
        line_width_pt = 2  # Width of the circle outline in points
        # Calculate circle position - centered on the shape

        # Add an oval shape (circle)
        circle = slide.shapes.add_shape(
            autoshape_type_id=1,  # MSO_SHAPE_OVAL (value 1)
            left=shape.left,
            top=shape.top,
            width=shape.width,
            height=shape.height,
        )

        # Set no fill for the circle (transparent inside)
        circle.fill.background()  # or circle.fill.solid() + set transparency

        # Set outline color and width
        circle.line.color.rgb = RGBColor(*line_color_rgb)
        circle.line.width = Pt(line_width_pt)

        self._set_alt_text(
            circle,
            {"parent_shape_id": shape.shape_id, "meta": True},
        )

    def _select_slice_from_df(self, df, integration):
        """
        Selects a specific slice from the DataFrame based on the integration settings.

        Args:
            df: A pandas DataFrame containing the data.
            integration: A LookerReference object containing the integration settings.
        Returns:
            The selected data slice (str or other type).
        """
        if integration.row is not None:
            row_slice = integration.row
        else:
            row_slice = 0

        row = df.iloc[row_slice]

        if integration.label is not None and integration.column is not None:
            logging.warning(
                f"Both label and column are set for integration {integration.id}. Defaulting to label and ignoring column."
            )
            r = row[integration.label]
        elif integration.label is not None:
            r = row[integration.label]
        elif integration.column is not None:
            r = row.iloc[integration.column]
        else:
            r = df
        return r

    def _replace_image_with_object(
        self, slide_index, shape_number, image_stream, integration
    ):
        slide = self.presentation.slides[slide_index]
        old_shape = next((s for s in slide.shapes if s.shape_id == shape_number), None)
        if old_shape is None:
            raise ValueError(f"Shape {shape_number} not found on slide {slide_index}.")
        if old_shape.shape_type != 13:  # picture
            raise ValueError("Selected shape is not an image.")

        left, top, width, height = (
            old_shape.left,
            old_shape.top,
            old_shape.width,
            old_shape.height,
        )
        slide.shapes._spTree.remove(old_shape._element)

        # --- calculate scaled size preserving aspect ratio ---
        img_bytes = image_stream.getvalue()
        image_stream.seek(0)
        with Image.open(BytesIO(img_bytes)) as im:
            img_w, img_h = im.size
        img_ratio = img_w / img_h
        shape_ratio = width / height

        if img_ratio > shape_ratio:
            new_width = width
            new_height = int(width / img_ratio)
        else:
            new_height = height
            new_width = int(height * img_ratio)

        # center within original box
        new_left = left + (width - new_width) / 2
        new_top = top + (height - new_height) / 2

        picture = slide.shapes.add_picture(
            BytesIO(img_bytes), new_left, new_top, width=new_width, height=new_height
        )
        self._set_alt_text(picture, integration)

    def _remove_shape(self, slide_index, shape_number):
        """
        Removes a shape from a PowerPoint slide.
        Args:
            prs: The Presentation object.
            slide_index: The index of the slide containing the shape.
            shape_index: The index of the shape to remove.
        """

        slide = self.presentation.slides[slide_index]
        shape_to_remove = None
        for shape in slide.shapes:
            if shape.shape_id == shape_number:
                shape_to_remove = shape

        if shape_to_remove is None:
            raise ValueError(
                f"Shape with number {shape_number} not found on slide {slide_index}."
            )

        # Remove the shape
        slide.shapes._spTree.remove(shape_to_remove._element)

    def _format_context_data(self, df) -> str:
        """
        Format a pandas DataFrame as a human-readable plain-text table for use
        as Gemini context.

        Args:
            df: A pandas DataFrame.

        Returns:
            str: A plain-text representation of the DataFrame.
        """
        return df.to_string(index=False)

    def _extract_slide_text_context(
        self, slide_number: int, exclude_shape_id: int
    ) -> str:
        """
        Return a plain-text extract of all text on a slide, excluding the target shape.

        Only shapes that have a text_frame are included.  Each shape is rendered as:
        ``[Shape name]: text content``

        Args:
            slide_number: The index of the slide.
            exclude_shape_id: The shape_id of the shape to exclude (the Gemini shape itself).

        Returns:
            str: Multi-line text representing the slide content.
        """
        slide = self.presentation.slides[slide_number]
        lines: list[str] = []
        for shape in slide.shapes:
            if shape.shape_id == exclude_shape_id:
                continue
            if hasattr(shape, "text_frame"):
                text = shape.text_frame.text.strip()
                if text:
                    name = getattr(shape, "name", f"shape_{shape.shape_id}")
                    lines.append(f"[{name}]: {text}")
        return "\n".join(lines)

    def _resolve_context_item(
        self,
        ctx: str,
        shape_number: int,
        slide_number: int,
        gemini_results: dict,
        current_text: str,
    ) -> tuple | None:
        """
        Resolve a single ``contexts`` entry to a ``(label, content)`` pair.

        Resolution rules (checked in order):

        1. ``"self"`` — the shape's own current text before synthesis.
        2. ``"slide_self"`` — text of all other shapes on the slide after Looker
           data has been rendered, with this shape excluded.
        3. Strings starting with ``"gemini_"`` — the output of the Gemini box
           whose ``gemini_id`` matches.  Returns ``None`` (with a warning) if
           that box has not been processed yet or has failed.
        4. Anything else — treated as a Looker meta-look ``meta_name``; resolved
           from ``self.data``.  Returns ``None`` (with a warning) if missing.

        Args:
            ctx: The context string from ``GeminiConfig.contexts``.
            shape_number: The shape_id of the current Gemini shape (for exclusion).
            slide_number: The slide index of the current Gemini shape.
            gemini_results: Dict of ``{gemini_id: synthesized_text}`` for already
                processed Gemini boxes.
            current_text: The shape's current text content.

        Returns:
            ``(label, content)`` on success, or ``None`` if the reference cannot
            be resolved (a warning is already logged).
        """
        if ctx == "self":
            return ("Current shape text", current_text)

        if ctx == "slide_self":
            content = self._extract_slide_text_context(slide_number, shape_number)
            return ("Current slide context (excluding this shape)", content)

        if ctx.startswith("gemini_"):
            result = gemini_results.get(ctx)
            if result is None:
                logging.warning(
                    f"No result available for Gemini context '{ctx}'. "
                    "The referenced Gemini shape may not have run yet or may have failed."
                )
                return None
            return (f"LLM report [{ctx}]", result)

        # Meta-look fallback
        raw = self.data.get(ctx)
        if raw is None:
            logging.warning(
                f"No data found for Gemini context '{ctx}'. "
                "Make sure a meta-look shape with that meta_name exists in the presentation."
            )
            return None
        try:
            df = self._make_df(raw)
            return (f"Data [{ctx}]", self._format_context_data(df))
        except Exception as e:
            logging.warning(f"Could not format context data for '{ctx}': {e}")
            return None

    def _sort_gemini_shapes_by_dependency(self) -> list:
        """
        Return ``self.gemini_shapes`` sorted in topological order.

        Any ``contexts`` entry that starts with ``"gemini_"`` and matches the
        ``gemini_id`` of another shape in the list is treated as a dependency —
        that shape must be processed first.

        Raises
        ------
        ValueError
            If a circular dependency is detected.
        """
        id_to_shape: dict[str, object] = {
            gs.integration.gemini_id: gs
            for gs in self.gemini_shapes
            if gs.integration.gemini_id
        }

        in_degree: dict[int, int] = {id(gs): 0 for gs in self.gemini_shapes}
        dependents: dict[int, list] = {id(gs): [] for gs in self.gemini_shapes}

        for gs in self.gemini_shapes:
            for ctx in gs.integration.contexts:
                if ctx.startswith("gemini_") and ctx in id_to_shape:
                    dep = id_to_shape[ctx]
                    dependents[id(dep)].append(gs)
                    in_degree[id(gs)] += 1

        queue: collections.deque = collections.deque(
            gs for gs in self.gemini_shapes if in_degree[id(gs)] == 0
        )
        ordered: list = []
        while queue:
            node = queue.popleft()
            ordered.append(node)
            for dependent in dependents[id(node)]:
                in_degree[id(dependent)] -= 1
                if in_degree[id(dependent)] == 0:
                    queue.append(dependent)

        if len(ordered) != len(self.gemini_shapes):
            raise ValueError(
                "Circular dependency detected in Gemini shape contexts. "
                "Ensure there are no cycles among gemini_id references."
            )

        return ordered

    def _process_gemini_shapes(self):
        """
        Process all shapes configured for Gemini LLM synthesis.

        Shapes are processed in dependency order (any shape whose ``gemini_id``
        appears in another shape's ``contexts`` is processed first).

        For each GeminiShape the ``contexts`` list is walked in order and each
        entry resolved via :meth:`_resolve_context_item` — which handles
        ``"self"``, ``"slide_self"``, Gemini box references (``gemini_``-prefix),
        and Looker meta-look names.  The assembled context string is then sent to
        the Gemini API together with the shape's current text and prompt.

        On error: the error message is written into the text box and a red outline
        is drawn around the shape (suppressed by ``--hide-errors``).
        """
        if not self.gemini_shapes:
            return

        if not gemini_module.is_available():
            logging.warning(
                "google-genai is not installed; Gemini synthesis shapes will be skipped. "
                "Install it with 'pip install looker_powerpoint[llm]' to enable LLM features."
            )
            return

        try:
            ordered_shapes = self._sort_gemini_shapes_by_dependency()
        except ValueError as e:
            logging.error(f"Cannot process Gemini shapes: {e}")
            return

        # Stores synthesized text keyed by gemini_id for chaining
        gemini_results: dict[str, str] = {}

        for gemini_shape in ordered_shapes:
            slide = self.presentation.slides[gemini_shape.slide_number]
            current_shape = None
            for shape in slide.shapes:
                if shape.shape_id == gemini_shape.shape_number:
                    current_shape = shape
                    break

            if current_shape is None:
                logging.error(
                    f"Could not find shape {gemini_shape.shape_number} on slide "
                    f"{gemini_shape.slide_number} for Gemini synthesis."
                )
                continue

            try:
                # Capture current text first (needed by "self" resolver and as
                # replacement-target context for synthesize())
                current_text = ""
                if hasattr(current_shape, "text_frame"):
                    current_text = current_shape.text_frame.text
                else:
                    logging.warning(
                        f"Shape {gemini_shape.shape_number} on slide "
                        f"{gemini_shape.slide_number} has no text_frame; "
                        "skipping Gemini synthesis."
                    )
                    continue

                # Resolve each context entry in order
                context_parts: list[str] = []
                for ctx in gemini_shape.integration.contexts:
                    resolved = self._resolve_context_item(
                        ctx,
                        gemini_shape.shape_number,
                        gemini_shape.slide_number,
                        gemini_results,
                        current_text,
                    )
                    if resolved is not None:
                        label, content = resolved
                        if content:
                            context_parts.append(f"{label}:\n{content}")

                context_data_str = "\n\n".join(context_parts)

                synthesized = gemini_module.synthesize(
                    prompt=gemini_shape.integration.prompt,
                    context_data_str=context_data_str,
                    current_text=current_text,
                    model_name=gemini_shape.integration.model,
                )

                update_text_frame_preserving_formatting(
                    current_shape.text_frame, synthesized
                )
                logging.debug(
                    f"Gemini synthesis applied to shape {gemini_shape.shape_number} "
                    f"on slide {gemini_shape.slide_number}."
                )

                # Store result for downstream shapes
                if gemini_shape.integration.gemini_id:
                    gemini_results[gemini_shape.integration.gemini_id] = synthesized

            except Exception as e:
                error_msg = str(e)
                logging.error(
                    f"Gemini synthesis failed for shape {gemini_shape.shape_number} "
                    f"on slide {gemini_shape.slide_number}: {error_msg}"
                )
                # Populate error message into text box
                try:
                    if hasattr(current_shape, "text_frame"):
                        update_text_frame_preserving_formatting(
                            current_shape.text_frame, error_msg
                        )
                except Exception:
                    pass
                # Draw red outline around the failed shape
                if not self.args.hide_errors:
                    self._mark_failure(slide, current_shape)

    def _make_df(self, result):
        """
        Create a pandas DataFrame from Looker data based on the integration settings.
        Categorizes and sorts columns into Dimensions -> Pivots -> Table Calcs.
        """
        data = json.loads(result)
        fields = data.get("metadata", {}).get("fields", {})

        # 1. Pull the injected sorts and pivots rules from the Look
        look_sorts = data.get("custom_sorts", [])
        look_pivots = data.get("custom_pivots", [])

        # Determine if the primary pivot is sorted descending
        pivot_descending = False
        if look_pivots:
            main_pivot = look_pivots[0]
            # Looker sorts look like: "view_name.date_dim desc 0"
            # Parse sort strings by tokens so we only match the exact field name,
            # not arbitrary substrings.
            for sort_str in look_sorts:
                parts = str(sort_str).split()
                if len(parts) < 2:
                    continue
                field_token = parts[0]
                # Find the first explicit direction token after the field
                direction_token = None
                for token in parts[1:]:
                    token_lower = token.lower()
                    if token_lower in ("asc", "desc"):
                        direction_token = token_lower
                        break
                if field_token == main_pivot and direction_token == "desc":
                    pivot_descending = True
                    break

        # Create DataFrame first to expose all dynamic column names
        df = pd.json_normalize(data.get("rows", [])).fillna("")
        actual_cols = list(df.columns)

        # 2. Extract base names from the metadata
        dim_bases = [f["name"] for f in fields.get("dimensions", [])]
        calc_bases = [f["name"] for f in fields.get("table_calculations", [])]
        measure_bases = [f["name"] for f in fields.get("measures", [])]

        dims = []
        calcs = []
        pivots_and_measures = []
        leftovers = []

        for col in actual_cols:
            if any(col == f"{d}.value" or col == d for d in dim_bases):
                dims.append(col)
            elif any(col == f"{c}.value" or col == c for c in calc_bases):
                calcs.append(col)
            elif "|FIELD|" in col or any(
                col == f"{m}.value" or col == m for m in measure_bases
            ):
                pivots_and_measures.append(col)
            else:
                leftovers.append(col)

        # 3. Apply Dimensions and Calcs sorting (Native query order)
        dims.sort(
            key=lambda x: next(
                (i for i, d in enumerate(dim_bases) if x == f"{d}.value" or x == d), 999
            )
        )
        calcs.sort(
            key=lambda x: next(
                (i for i, c in enumerate(calc_bases) if x == f"{c}.value" or x == c),
                999,
            )
        )

        # 4. Apply Looker's strict Pivot Sorting Rules
        def parse_pivot_col(col):
            # Break down "measure_name|FIELD|2025-03-03.value"
            base = col.replace(".value", "")
            if "|FIELD|" in base:
                measure, pivot_val = base.split("|FIELD|", 1)
                return pivot_val, measure
            return "", base

        # Get unique pivot values preserving their order of first appearance in the data
        # (which reflects Looker's native ordering). Avoid lexicographic sorting so that
        # numeric-like values ("2", "10") and non-ISO dates aren't misordered.
        seen_pivots: set = set()
        unique_pivots: list = []
        for c in pivots_and_measures:
            pv = parse_pivot_col(c)[0]
            if pv not in seen_pivots:
                seen_pivots.add(pv)
                unique_pivots.append(pv)
        if pivot_descending:
            unique_pivots.reverse()

        pivot_order_map = {val: i for i, val in enumerate(unique_pivots)}
        measure_order_map = {m: i for i, m in enumerate(measure_bases)}

        # Sort first by the properly sequenced pivot value, then by the measure's native query order
        pivots_and_measures.sort(
            key=lambda x: (
                pivot_order_map.get(parse_pivot_col(x)[0], 999),
                measure_order_map.get(parse_pivot_col(x)[1], 999),
            )
        )

        # 5. Re-index and Rename
        ordered_cols = dims + pivots_and_measures + calcs + leftovers
        df = df[ordered_cols]

        all_fields = (
            fields.get("dimensions", [])
            + fields.get("measures", [])
            + fields.get("table_calculations", [])
        )
        mappy = {
            f"{item['name']}.value": item.get("field_group_variant", item["name"])
            .strip()
            .lower()
            .replace(" ", "_")
            for item in all_fields
        }
        df.rename(columns=mappy, inplace=True)

        return df

    def _build_metadata_object(self):
        """
        Build metadata object for the presentation.
        """
        metadata_rows = []
        looks = set()
        for looker_shape in self.looker_shapes:
            if looker_shape.integration.id not in looks:
                looks.add(looker_shape.integration.id)
                metadata_rows.append(
                    {
                        "looks": {
                            "value": f"{os.environ.get('LOOKERSDK_BASE_URL')}looks/{looker_shape.integration.id}"
                        }
                    }
                )
        metadata_object = {
            "metadata": {"fields": {"dimensions": [{"name": "looks"}]}},
            "rows": metadata_rows,
        }
        self.data["metadata_shapes"] = json.dumps(metadata_object)

    async def get_queries(self):
        """
        asynchronously fetch a list of look references
        """
        logging.info(
            f"Running Looker queries... {len(self.looker_shapes)} queries to run."
        )
        tasks = [
            self.client._async_write_queries(
                shape.shape_id, self.args.filter, **dict(shape.integration)
            )
            for shape in self.looker_shapes
        ]

        # Run all tasks concurrently and gather the results
        results = await asyncio.gather(*tasks)
        for r in results:
            self.data.update(r)

    def _test_str_to_int(self, s):
        try:
            int(s)
            return True
        except ValueError:
            return False

    def run(self, **kwargs):
        """
        Main method to run the CLI application.
        """
        self.args = self.parser.parse_args()
        self._setup_logging()
        self._pick_file()
        self._init_looker()

        references = self.get_alt_text(self.file_path)
        if not references:
            logging.error(
                "No shapes with id found in the presentation. Add a 'id' : '<look_id>' to the alternative text of a shape to load data into the shape."
            )
            return

        for ref in references:
            integration = ref.get("integration", {})
            # Try to parse as a Gemini shape first (type: gemini discriminator)
            if isinstance(integration, dict) and integration.get("type") == "gemini":
                try:
                    gemini_shape = GeminiShape.model_validate(ref)
                    if gemini_shape.shape_type not in (
                        "TEXT_BOX",
                        "TITLE",
                        "AUTO_SHAPE",
                    ):
                        logging.warning(
                            f"Gemini synthesis config found on shape "
                            f"{gemini_shape.shape_id} (type: {gemini_shape.shape_type}). "
                            "Gemini synthesis only works for text boxes (TEXT_BOX, TITLE, "
                            "AUTO_SHAPE). This shape will be skipped."
                        )
                        continue
                    self.gemini_shapes.append(gemini_shape)
                except ValidationError as e:
                    logging.debug(
                        f"Could not parse Gemini config in shape {ref.get('shape_id', '?')}: {e}"
                    )
                continue

            # Otherwise try to parse as a regular Looker shape
            try:
                self.relevant_shapes.append(LookerShape.model_validate(ref))
            except ValidationError as e:
                logging.debug(
                    f"Could not parse the alternate text in slide {ref['shape_id'].split(',')[0]}, shape {ref['shape_id'].split(',')[1]}: {e}"
                )
                continue

        self.looker_shapes = [
            s
            for s in self.relevant_shapes
            if s.integration.id_type == "look"
            and self._test_str_to_int(s.integration.id)
        ]

        self._build_metadata_object()

        asyncio.run(self.get_queries())

        for looker_shape in self.relevant_shapes:
            if looker_shape.integration.meta:
                if not self.args.self:
                    self._remove_shape(
                        looker_shape.slide_number,
                        looker_shape.shape_number,
                    )

            else:
                result = self.data.get(looker_shape.shape_id)
                if result is None:
                    result = self.data.get(looker_shape.integration.id)

                try:
                    if looker_shape.shape_type == "PICTURE":
                        if looker_shape.integration.result_format in ("jpg", "png"):
                            image_stream = BytesIO(result)
                        else:
                            df = self._make_df(result)
                            url = self._select_slice_from_df(
                                df, looker_shape.integration
                            )

                            response = requests.get(url)
                            response.raise_for_status()
                            image_stream = io.BytesIO(response.content)

                        logging.debug(
                            f"Replacing image for shape {looker_shape.shape_number} on slide {looker_shape.slide_number}..."
                        )

                        self._replace_image_with_object(
                            looker_shape.slide_number,
                            looker_shape.shape_number,
                            image_stream,
                            looker_shape.original_integration,
                        )

                    elif looker_shape.shape_type in [
                        "CHART",
                        "TABLE",
                        "TEXT_BOX",
                        "TITLE",
                        "AUTO_SHAPE",
                    ]:
                        slide = self.presentation.slides[looker_shape.slide_number]
                        for shape in slide.shapes:
                            if shape.shape_id == looker_shape.shape_number:
                                current_shape = shape
                        df = self._make_df(result)

                        if looker_shape.shape_type == "TABLE":
                            logging.debug(
                                f"Updating table for shape {looker_shape.shape_number} on slide {looker_shape.slide_number}..."
                            )
                            self._fill_table(
                                current_shape.table,
                                df,
                                looker_shape.integration.headers,
                            )

                        elif looker_shape.shape_type in [
                            "TEXT_BOX",
                            "TITLE",
                            "AUTO_SHAPE",
                        ]:
                            logging.debug(
                                f"Updating text for shape {looker_shape.shape_number} on slide {looker_shape.slide_number}..."
                            )

                            try:
                                text_to_insert = self._select_slice_from_df(
                                    df, looker_shape.integration
                                )
                            except Exception as e:
                                text_to_insert = df.to_string(index=False, header=False)
                                logging.debug(
                                    f"inserting whole text for shape {looker_shape.shape_number} on slide {looker_shape.slide_number}: {e}"
                                )
                            current_shape = process_text_field(
                                current_shape,
                                text_to_insert,
                                df,
                            )
                            # add_text_with_numbered_links(current_shape.text_frame, str(text_to_insert))

                        elif looker_shape.shape_type == "CHART":
                            chart_data = CategoryChartData()
                            chart_data.categories = df.iloc[
                                :, 0
                            ].tolist()  # Assuming the first column contains categories
                            chart = current_shape.chart
                            existing_chart_data = chart.plots[0].series
                            logging.debug(
                                f"Existing chart series: {[s.name for s in existing_chart_data]}"
                            )

                            if looker_shape.integration.headers:
                                for series_name in df.columns[1:]:
                                    try:
                                        match = (
                                            re.search(
                                                r"^[^\.]*\.[^\.]*\.(.*)\.value$",
                                                series_name,
                                            )
                                            .group(1)
                                            .replace(".", " - ")
                                            .strip()
                                            .replace("|FIELD|", " ")
                                        )
                                    except AttributeError as e:
                                        logging.debug(
                                            f"Could not parse series name {series_name}, setting name to {series_name}"
                                        )
                                        match = series_name
                                    chart_data.add_series(match, df[series_name])
                            else:
                                if len(df.columns[1:]) != len(existing_chart_data):
                                    logging.warning(
                                        f"{looker_shape.shape_id}. Missing headers! Number of series ({len(df.columns[1:])}) does not match number of existing chart series ({len(existing_chart_data)}). Perhaps you need to enable headers in the integration settings?"
                                    )
                                for series_name, series in zip(
                                    df.columns[1:], existing_chart_data
                                ):
                                    chart_data.add_series(series.name, df[series_name])

                            chart.replace_data(chart_data)
                            if looker_shape.integration.show_latest_chart_label:
                                for plot in chart.plots:
                                    s = 0
                                    for series in plot.series:
                                        series_has_label = False
                                        index = 0
                                        for i, v in zip(
                                            series.points, df.iloc[:, s + 1]
                                        ):
                                            if i.data_label._dLbl is not None:
                                                series_has_label = True
                                                logging.debug(
                                                    f"Series {series.name} has data labels."
                                                )
                                            if v is not None and v != "":
                                                logging.debug(
                                                    f"Value for point {index} in series {series.name}: {v}"
                                                )
                                                index += 1
                                        if series_has_label is True:
                                            new_index = 0
                                            for point in series.points:
                                                new_index += 1
                                                if new_index == index:
                                                    logging.debug(
                                                        f"Showing data label for point {new_index} in series {series.name}."
                                                    )
                                                    point.data_label.text_frame.text = (
                                                        ""
                                                    )
                                                    point.data_label.has_text_frame = (
                                                        False
                                                    )
                                                else:
                                                    point.data_label.text_frame.text = (
                                                        ""
                                                    )
                                                    point.data_label.has_text_frame = (
                                                        True
                                                    )
                                        s += 1

                    else:
                        logging.warning(
                            f"unknown shape type {looker_shape.shape_type} for shape {looker_shape.shape_number} on slide {looker_shape.slide_number}."
                        )
                        continue

                except Exception as e:
                    logging.error(f"Error processing reference {looker_shape}: {e}")
                    # import traceback
                    # traceback.print_exc()  # Prints the full traceback

                    if not self.args.hide_errors:
                        slide = self.presentation.slides[looker_shape.slide_number]
                        for shape in slide.shapes:
                            if shape.shape_id == looker_shape.shape_number:
                                self._mark_failure(slide, shape)

        # Process Gemini synthesis shapes
        self._process_gemini_shapes()

        if self.args.self:
            self.destination = self.file_path
        else:
            if not self.args.output_dir.endswith("/"):
                self.args.output_dir += "/"
            self.destination = (
                self.args.output_dir
                + os.path.basename(self.file_path).removesuffix(".pptx")
                + datetime.datetime.now().strftime("_%Y%m%d_%H%M%S.pptx")
            )
        if not os.path.exists(self.args.output_dir) and not self.args.self:
            os.makedirs(self.args.output_dir)

        self.presentation.save(self.destination)

        if not self.args.quiet:
            try:
                os.startfile(self.destination)
                logging.info(f"Opened {self.destination} in PowerPoint.")
            except Exception as e:
                try:
                    subprocess.Popen(["open", self.destination])  # For macOS
                    logging.info(f"Opened {self.destination} in PowerPoint.")
                except Exception as e:
                    logging.error(f"Failed to open the PowerPoint file: {e}")
                    logging.info(f"You can find the file at {self.destination}.")


def main():
    cli = Cli()
    cli.run()


if __name__ == "__main__":
    main()
