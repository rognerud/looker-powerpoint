"""Unit tests for the pptx test fixtures.

Tests in this module validate assumptions about the pptx files stored in
``test/pptx/``.  They exercise:

* :func:`~looker_powerpoint.tools.find_alt_text.extract_alt_text` — low-level
  per-shape YAML extraction.
* :func:`~looker_powerpoint.tools.find_alt_text.get_presentation_objects_with_descriptions`
  — full-presentation parse that returns a list of shape dicts.

No live Looker API calls are made; all tests operate on the local pptx fixtures.
"""

import os
from types import SimpleNamespace

import pytest
from pptx import Presentation

from looker_powerpoint.tools.find_alt_text import (
    cleanse_alt_text,
    extract_alt_text,
    get_presentation_objects_with_descriptions,
)

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

PPTX_DIR = os.path.join(os.path.dirname(__file__), "pptx")
TABLE7X7_PATH = os.path.join(PPTX_DIR, "table7x7.pptx")

# EMUs (English Metric Units) per pixel, as used by python-pptx
EMU_PER_PIXEL = 9525


# ---------------------------------------------------------------------------
# TestTable7x7Pptx — structural assumptions about table7x7.pptx
# ---------------------------------------------------------------------------


class TestTable7x7Pptx:
    """Tests that validate the assumptions documented in test/pptx/table7x7.md."""

    def test_file_exists(self):
        """The fixture file must be present on disk."""
        assert os.path.isfile(TABLE7X7_PATH), f"Missing fixture: {TABLE7X7_PATH}"

    def test_presentation_has_one_slide(self):
        """table7x7.pptx contains exactly one slide."""
        prs = Presentation(TABLE7X7_PATH)
        assert len(prs.slides) == 1

    def test_slide_has_one_shape(self):
        """The single slide contains exactly one shape."""
        prs = Presentation(TABLE7X7_PATH)
        assert len(prs.slides[0].shapes) == 1

    def test_shape_is_table(self):
        """The shape on slide 0 must be a TABLE."""
        prs = Presentation(TABLE7X7_PATH)
        shape = prs.slides[0].shapes[0]
        assert shape.shape_type.name == "TABLE"

    def test_table_has_seven_rows(self):
        """The table must have 7 rows."""
        prs = Presentation(TABLE7X7_PATH)
        shape = prs.slides[0].shapes[0]
        assert shape.has_table
        assert len(shape.table.rows) == 7

    def test_table_has_seven_columns(self):
        """The table must have 7 columns."""
        prs = Presentation(TABLE7X7_PATH)
        shape = prs.slides[0].shapes[0]
        assert shape.has_table
        assert len(shape.table.columns) == 7

    def test_shape_alt_text_parses_to_dict(self):
        """extract_alt_text returns a dict (parsed YAML), not None."""
        prs = Presentation(TABLE7X7_PATH)
        shape = prs.slides[0].shapes[0]
        result = extract_alt_text(shape)
        assert isinstance(result, dict)

    def test_shape_alt_text_has_id_1(self):
        """The YAML alt text sets ``id: 1`` as documented in table7x7.md."""
        prs = Presentation(TABLE7X7_PATH)
        shape = prs.slides[0].shapes[0]
        result = extract_alt_text(shape)
        assert result == {"id": 1}

    def test_shape_dimensions_in_pixels(self):
        """Shape width and height match the expected pixel values."""
        prs = Presentation(TABLE7X7_PATH)
        shape = prs.slides[0].shapes[0]
        emu_to_px = lambda emu: round(emu / EMU_PER_PIXEL)
        assert emu_to_px(shape.width) == 853
        assert emu_to_px(shape.height) == 273


# ---------------------------------------------------------------------------
# TestGetPresentationObjects — get_presentation_objects_with_descriptions
# ---------------------------------------------------------------------------


class TestGetPresentationObjects:
    """Tests for get_presentation_objects_with_descriptions using table7x7.pptx."""

    @pytest.fixture(scope="class")
    def objects(self):
        return get_presentation_objects_with_descriptions(TABLE7X7_PATH)

    def test_returns_one_object(self, objects):
        """Exactly one shape with alt text is present in the presentation."""
        assert len(objects) == 1

    def test_shape_type_is_table(self, objects):
        """The extracted shape type is TABLE."""
        assert objects[0]["shape_type"] == "TABLE"

    def test_slide_number_is_zero(self, objects):
        """The shape lives on slide index 0."""
        assert objects[0]["slide_number"] == 0

    def test_shape_id_format(self, objects):
        """shape_id is formatted as '<slide_index>,<shape_id>'."""
        assert objects[0]["shape_id"] == "0,4"

    def test_shape_number(self, objects):
        """shape_number matches the pptx shape id attribute."""
        assert objects[0]["shape_number"] == 4

    def test_integration_id(self, objects):
        """The integration dict contains the parsed YAML id value."""
        assert objects[0]["integration"] == {"id": 1}

    def test_shape_width_pixels(self, objects):
        """shape_width is the expected pixel value."""
        assert objects[0]["shape_width"] == 853

    def test_shape_height_pixels(self, objects):
        """shape_height is the expected pixel value."""
        assert objects[0]["shape_height"] == 273


# ---------------------------------------------------------------------------
# TestExtractAltText — edge cases for extract_alt_text
# ---------------------------------------------------------------------------


class TestExtractAltText:
    """Edge-case tests for extract_alt_text."""

    def test_shape_without_alt_text_returns_none(self):
        """A shape that carries no alt-text description returns None."""
        prs = Presentation(TABLE7X7_PATH)
        shape = prs.slides[0].shapes[0]
        # Remove the descr attribute from the cNvPr element to simulate a
        # shape without alt text, then confirm the function returns None.
        from lxml import etree

        xml_elem = etree.fromstring(shape.element.xml)
        NS = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
        for path in [
            ".//p:nvSpPr/p:cNvPr",
            ".//p:nvPicPr/p:cNvPr",
            ".//p:nvGraphicFramePr/p:cNvPr",
        ]:
            for el in xml_elem.xpath(path, namespaces=NS):
                if "descr" in el.attrib:
                    del el.attrib["descr"]
        # Re-create a minimal mock shape whose .element.xml returns the stripped XML
        fake_shape = SimpleNamespace(
            element=SimpleNamespace(xml=etree.tostring(xml_elem, encoding="unicode"))
        )

        assert extract_alt_text(fake_shape) is None

    def test_invalid_path_returns_empty_list(self):
        """An invalid file path to get_presentation_objects_with_descriptions returns []."""
        result = get_presentation_objects_with_descriptions("/nonexistent/path.pptx")
        assert result == []


# ---------------------------------------------------------------------------
# TestCleanseAltText — unit tests for cleanse_alt_text
# ---------------------------------------------------------------------------


class TestCleanseAltText:
    """Tests for the cleanse_alt_text helper that normalises typographic quotes."""

    def test_left_double_quotation_mark(self):
        assert cleanse_alt_text("\u201cvalue\u201d") == '"value"'

    def test_right_double_quotation_mark(self):
        assert cleanse_alt_text("key: \u201cval\u201d") == 'key: "val"'

    def test_double_low9_quotation_mark(self):
        assert cleanse_alt_text("\u201eval\u201d") == '"val"'

    def test_left_single_quotation_mark(self):
        assert cleanse_alt_text("\u2018hello\u2019") == "'hello'"

    def test_right_single_quotation_mark(self):
        assert cleanse_alt_text("it\u2019s") == "it's"

    def test_single_low9_quotation_mark(self):
        assert cleanse_alt_text("\u201atest\u2019") == "'test'"

    def test_prime_replaced_with_apostrophe(self):
        assert cleanse_alt_text("O\u2032clock") == "O'clock"

    def test_double_prime_replaced_with_double_quote(self):
        assert cleanse_alt_text("value\u2033") == 'value"'

    def test_grave_accent_replaced(self):
        assert cleanse_alt_text("\u0060key\u0060") == "'key'"

    def test_acute_accent_replaced(self):
        assert cleanse_alt_text("\u00b4key\u00b4") == "'key'"

    def test_angle_double_quotes_replaced(self):
        assert cleanse_alt_text("\u00abvalue\u00bb") == '"value"'

    def test_plain_text_unchanged(self):
        assert cleanse_alt_text("id: 42") == "id: 42"

    def test_mixed_fancy_and_plain_quotes(self):
        result = cleanse_alt_text("id: \u201c42\u201d\nlabel: \u2018foo\u2019")
        assert result == "id: \"42\"\nlabel: 'foo'"

    def test_yaml_with_smart_quotes_parses_correctly(self):
        """Smart-quote YAML should parse to the same dict as plain-quote YAML."""
        import yaml

        fancy = "id: \u201c42\u201d"
        plain = 'id: "42"'
        assert yaml.safe_load(cleanse_alt_text(fancy)) == yaml.safe_load(plain)

    def test_extract_alt_text_with_smart_quotes(self):
        """extract_alt_text tolerates smart-quote YAML embedded in a shape."""
        # Build a minimal shape XML that uses curly quotes around the id value
        xml_str = (
            '<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
            ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            "<p:nvSpPr>"
            '<p:cNvPr id="1" name="TextBox 1" descr="id: \u201c99\u201d"/>'
            "</p:nvSpPr>"
            "</p:sp>"
        )
        fake_shape = SimpleNamespace(element=SimpleNamespace(xml=xml_str))
        result = extract_alt_text(fake_shape)
        assert result == {"id": "99"}

    def test_empty_string_unchanged(self):
        assert cleanse_alt_text("") == ""
