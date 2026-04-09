"""Integration test for the full CLI pipeline.

Exercises the end-to-end flow:
  read pptx  →  mock Looker API  →  run cli.run()  →  verify filled output pptx
"""

import argparse
import json
import os
from unittest.mock import AsyncMock, MagicMock, patch

from pptx import Presentation

from looker_powerpoint.cli import Cli

# ── constants ─────────────────────────────────────────────────────────────────

PPTX_PATH = os.path.join(os.path.dirname(__file__), "pptx", "table7x7.pptx")
# Shape ID returned by get_presentation_objects_with_descriptions for the
# single TABLE shape in table7x7.pptx (slide index 0, pptx shape_id 4).
TABLE_SHAPE_ID = "0,4"


# ── helpers ───────────────────────────────────────────────────────────────────


def _json_bi(dimensions, measures, table_calculations, rows):
    """Build a minimal json_bi-format payload (same structure as the helper in test_cli.py)."""

    def _field(name):
        return {"name": name, "field_group_variant": name.split(".")[-1]}

    return json.dumps(
        {
            "metadata": {
                "fields": {
                    "dimensions": [_field(d) for d in dimensions],
                    "measures": [_field(m) for m in measures],
                    "table_calculations": [
                        _field(t) for t in (table_calculations or [])
                    ],
                }
            },
            "rows": rows,
            "custom_sorts": [],
            "custom_pivots": [],
        }
    )


def _make_args(pptx_path, output_dir):
    """Return an `argparse.Namespace` that matches every attribute read by `Cli.run`."""
    ns = argparse.Namespace(
        file_path=pptx_path,
        output_dir=output_dir,
        add_links=False,
        hide_errors=True,
        parse_date_syntax_in_filename=False,
        quiet=True,
        filter=None,
        debug_queries=False,
        verbose=0,
    )
    # "self" is not a Python keyword, but using it as a kwarg looks odd; setattr is cleaner.
    setattr(ns, "self", False)
    return ns


# ── integration test ──────────────────────────────────────────────────────────


class TestIntegration:
    """End-to-end integration tests: pptx → mocked Looker → filled pptx."""

    def test_run_fills_table_from_mocked_looker(self, tmp_path):
        """Full pipeline: parse table7x7.pptx, return mock Looker data, verify output table."""
        mock_result = _json_bi(
            dimensions=["orders.date", "orders.status"],
            measures=["orders.revenue", "orders.count"],
            table_calculations=[],
            rows=[
                {
                    "orders.date.value": "2024-01-01",
                    "orders.status.value": "complete",
                    "orders.revenue.value": "100",
                    "orders.count.value": "5",
                },
                {
                    "orders.date.value": "2024-01-02",
                    "orders.status.value": "pending",
                    "orders.revenue.value": "200",
                    "orders.count.value": "10",
                },
            ],
        )

        args = _make_args(PPTX_PATH, str(tmp_path))

        cli = Cli()
        # Override parse_args so pytest's own argv does not clash with the CLI parser.
        cli.parser.parse_args = lambda: args

        mock_client = MagicMock()
        mock_client._async_write_queries = AsyncMock(
            return_value={TABLE_SHAPE_ID: mock_result}
        )

        with patch("looker_powerpoint.cli.LookerClient", return_value=mock_client):
            cli.run()

        # ── verify output file ────────────────────────────────────────────────
        output_files = list(tmp_path.glob("*.pptx"))
        assert len(output_files) == 1, "Expected exactly one output pptx file"

        prs = Presentation(str(output_files[0]))
        table = next(s.table for s in prs.slides[0].shapes if s.has_table)

        # Header row reflects the field_group_variant names (part after last dot).
        assert table.cell(0, 0).text == "date"
        assert table.cell(0, 1).text == "status"
        assert table.cell(0, 2).text == "revenue"
        assert table.cell(0, 3).text == "count"

        # First data row
        assert table.cell(1, 0).text == "2024-01-01"
        assert table.cell(1, 1).text == "complete"
        assert table.cell(1, 2).text == "100"
        assert table.cell(1, 3).text == "5"

        # Second data row
        assert table.cell(2, 0).text == "2024-01-02"
        assert table.cell(2, 1).text == "pending"
        assert table.cell(2, 2).text == "200"
        assert table.cell(2, 3).text == "10"
