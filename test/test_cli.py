import json
import pandas as pd
import pytest
from unittest.mock import patch
from pptx import Presentation
from pptx.util import Inches
from looker_powerpoint.cli import Cli
from looker_powerpoint.models import LookerReference, LookerShape


def test_default_output_dir():
    """Test that the default output directory is 'output'."""
    with patch("os.getenv", return_value="dummy_value"):
        cli = Cli()
        args = cli.parser.parse_args()
        assert args.output_dir == "output"


def _make_cli():
    """Create a Cli instance with os.getenv stubbed out so no real environment is needed."""
    with patch("os.getenv", return_value="dummy_value"):
        return Cli()


def _field(name):
    """Build a field dict with a field_group_variant derived from the short name."""
    return {"name": name, "field_group_variant": name.split(".")[-1]}


def _make_result(dimensions, measures, table_calculations, rows, custom_sorts=None, custom_pivots=None):
    """Build a json_bi-style result string for _make_df.

    Each entry in *dimensions*, *measures*, and *table_calculations* may be either a
    plain field name string or a dict with ``name``/``field_group_variant`` keys.
    Plain strings are automatically expanded with a ``field_group_variant`` equal to
    the portion of the name after the last dot (e.g. ``"view.date"`` → ``"date"``).
    """
    def _normalise(fields):
        return [_field(f) if isinstance(f, str) else f for f in fields]

    return json.dumps({
        "metadata": {
            "fields": {
                "dimensions": _normalise(dimensions),
                "measures": _normalise(measures),
                "table_calculations": _normalise(table_calculations or []),
            }
        },
        "rows": rows,
        "custom_sorts": custom_sorts or [],
        "custom_pivots": custom_pivots or [],
    })


class TestMakeDf:
    """Tests for Cli._make_df column ordering logic."""

    def test_no_pivots_column_order(self):
        """Dimensions come first, then measures, then table calcs — no pivots."""
        cli = _make_cli()
        result = _make_result(
            dimensions=["view.date", "view.name"],
            measures=["view.revenue", "view.cost"],
            table_calculations=["calc1"],
            rows=[
                {
                    "view.date.value": "2024-01-01",
                    "view.name.value": "A",
                    "view.revenue.value": 100,
                    "view.cost.value": 50,
                    "calc1.value": 2.0,
                }
            ],
        )
        df = cli._make_df(result)
        cols = list(df.columns)
        # All five columns must be present (renamed via field_group_variant)
        assert len(cols) == 5
        # Dimensions first
        dim_positions = [cols.index("date"), cols.index("name")]
        # Measures next
        measure_positions = [cols.index("revenue"), cols.index("cost")]
        # Calc last
        calc_position = cols.index("calc1")
        assert max(dim_positions) < min(measure_positions), "Dims must precede measures"
        assert max(measure_positions) < calc_position, "Measures must precede table calcs"

    def test_no_pivots_preserves_native_dimension_order(self):
        """Dimension order follows the metadata field order, not data column order."""
        cli = _make_cli()
        result = _make_result(
            dimensions=["view.z_dim", "view.a_dim"],
            measures=["view.measure"],
            table_calculations=[],
            rows=[
                {
                    # a_dim appears first in the data dict on purpose
                    "view.a_dim.value": "foo",
                    "view.z_dim.value": "bar",
                    "view.measure.value": 1,
                }
            ],
        )
        df = cli._make_df(result)
        cols = list(df.columns)
        # z_dim was listed first in metadata → must appear before a_dim
        assert cols.index("z_dim") < cols.index("a_dim")

    def test_pivoted_measures_ascending(self):
        """Pivot columns are kept in their Looker-native (ascending) appearance order."""
        cli = _make_cli()
        # Looker returns pivot values in month order: Jan, Feb, Mar
        result = _make_result(
            dimensions=["view.region"],
            measures=["view.revenue"],
            table_calculations=[],
            rows=[
                {
                    "view.region.value": "North",
                    "view.revenue|FIELD|2024-01.value": 10,
                    "view.revenue|FIELD|2024-02.value": 20,
                    "view.revenue|FIELD|2024-03.value": 30,
                }
            ],
            custom_pivots=["view.month"],
            custom_sorts=["view.month asc 0"],
        )
        df = cli._make_df(result)
        pivot_cols = [c for c in df.columns if "revenue" in c]
        # The pivot order should match Looker's data order: Jan → Feb → Mar
        assert pivot_cols == [
            "view.revenue|FIELD|2024-01.value",
            "view.revenue|FIELD|2024-02.value",
            "view.revenue|FIELD|2024-03.value",
        ]

    def test_pivoted_measures_descending(self):
        """When the primary pivot is sorted desc, pivot column order is reversed."""
        cli = _make_cli()
        # Looker returns pivot values in month order: Jan, Feb, Mar
        # but sorts descending → we expect Mar, Feb, Jan
        result = _make_result(
            dimensions=["view.region"],
            measures=["view.revenue"],
            table_calculations=[],
            rows=[
                {
                    "view.region.value": "North",
                    "view.revenue|FIELD|2024-01.value": 10,
                    "view.revenue|FIELD|2024-02.value": 20,
                    "view.revenue|FIELD|2024-03.value": 30,
                }
            ],
            custom_pivots=["view.month"],
            custom_sorts=["view.month desc 0"],
        )
        df = cli._make_df(result)
        pivot_cols = [c for c in df.columns if "revenue" in c]
        # Descending → reversed from appearance order
        assert pivot_cols == [
            "view.revenue|FIELD|2024-03.value",
            "view.revenue|FIELD|2024-02.value",
            "view.revenue|FIELD|2024-01.value",
        ]

    def test_numeric_pivot_values_use_appearance_order_not_lexicographic(self):
        """
        Numeric-like pivot values ("2", "10") must be ordered by their appearance in
        Looker data, not lexicographically (which would give "10" before "2").
        """
        cli = _make_cli()
        # Looker returns data in numeric order: week 2, week 10
        result = _make_result(
            dimensions=["view.region"],
            measures=["view.revenue"],
            table_calculations=[],
            rows=[
                {
                    "view.region.value": "East",
                    "view.revenue|FIELD|2.value": 200,
                    "view.revenue|FIELD|10.value": 1000,
                }
            ],
            custom_pivots=["view.week"],
            custom_sorts=["view.week asc 0"],
        )
        df = cli._make_df(result)
        pivot_cols = [c for c in df.columns if "revenue" in c]
        # Should follow Looker's data order: 2 then 10
        # (lexicographic sort would give "10" before "2", which is wrong)
        assert pivot_cols == [
            "view.revenue|FIELD|2.value",
            "view.revenue|FIELD|10.value",
        ]

    def test_pivot_sort_only_reverses_on_exact_field_match(self):
        """
        A sort field that is a *substring* of the pivot field name must NOT trigger
        pivot_descending — only an exact token match should.
        """
        cli = _make_cli()
        # The sort field "view.month" is a substring of "view.month_long",
        # but only "view.month_long" is the actual pivot.
        result = _make_result(
            dimensions=["view.region"],
            measures=["view.revenue"],
            table_calculations=[],
            rows=[
                {
                    "view.region.value": "West",
                    "view.revenue|FIELD|Jan.value": 1,
                    "view.revenue|FIELD|Feb.value": 2,
                }
            ],
            custom_pivots=["view.month_long"],
            # "view.month" desc should NOT affect ordering of "view.month_long" pivot
            custom_sorts=["view.month desc 0"],
        )
        df = cli._make_df(result)
        pivot_cols = [c for c in df.columns if "revenue" in c]
        # No pivot_descending should be triggered → appearance order preserved
        assert pivot_cols == [
            "view.revenue|FIELD|Jan.value",
            "view.revenue|FIELD|Feb.value",
        ]


# ---------------------------------------------------------------------------
# Parser default / flag tests
# ---------------------------------------------------------------------------

class TestParser:
    """Tests for Cli._init_argparser argument defaults and flag overrides."""

    def test_default_output_dir(self):
        """Default output directory is 'output'."""
        cli = _make_cli()
        args = cli.parser.parse_args([])
        assert args.output_dir == "output"

    def test_default_file_path(self):
        """Default file path is None when not supplied."""
        cli = _make_cli()
        args = cli.parser.parse_args([])
        assert args.file_path is None

    def test_default_add_links(self):
        """--add-links defaults to False."""
        cli = _make_cli()
        args = cli.parser.parse_args([])
        assert args.add_links is False

    def test_default_hide_errors(self):
        """--hide-errors defaults to False."""
        cli = _make_cli()
        args = cli.parser.parse_args([])
        assert args.hide_errors is False

    def test_default_parse_date_syntax_in_filename(self):
        """--parse-date-syntax-in-filename defaults to True."""
        cli = _make_cli()
        args = cli.parser.parse_args([])
        assert args.parse_date_syntax_in_filename is True

    def test_default_self(self):
        """--self defaults to False."""
        cli = _make_cli()
        args = cli.parser.parse_args([])
        assert args.self is False

    def test_default_quiet(self):
        """--quiet defaults to False."""
        cli = _make_cli()
        args = cli.parser.parse_args([])
        assert args.quiet is False

    def test_default_filter(self):
        """--filter defaults to None."""
        cli = _make_cli()
        args = cli.parser.parse_args([])
        assert args.filter is None

    def test_default_debug_queries(self):
        """--debug-queries defaults to False."""
        cli = _make_cli()
        args = cli.parser.parse_args([])
        assert args.debug_queries is False

    def test_default_verbose(self):
        """-v / --verbose defaults to 0."""
        cli = _make_cli()
        args = cli.parser.parse_args([])
        assert args.verbose == 0

    def test_verbose_single_flag(self):
        """-v increments verbose to 1."""
        cli = _make_cli()
        args = cli.parser.parse_args(["-v"])
        assert args.verbose == 1

    def test_verbose_double_flag(self):
        """-vv increments verbose to 2."""
        cli = _make_cli()
        args = cli.parser.parse_args(["-vv"])
        assert args.verbose == 2

    def test_output_dir_long_flag(self):
        """--output-dir sets the output directory."""
        cli = _make_cli()
        args = cli.parser.parse_args(["--output-dir", "my_output"])
        assert args.output_dir == "my_output"

    def test_output_dir_short_flag(self):
        """-o is the short alias for --output-dir."""
        cli = _make_cli()
        args = cli.parser.parse_args(["-o", "my_output"])
        assert args.output_dir == "my_output"

    def test_file_path_long_flag(self):
        """--file-path sets the file path."""
        cli = _make_cli()
        args = cli.parser.parse_args(["--file-path", "test.pptx"])
        assert args.file_path == "test.pptx"

    def test_file_path_short_flag(self):
        """-f is the short alias for --file-path."""
        cli = _make_cli()
        args = cli.parser.parse_args(["-f", "test.pptx"])
        assert args.file_path == "test.pptx"

    def test_add_links_flag(self):
        """--add-links sets add_links to True."""
        cli = _make_cli()
        args = cli.parser.parse_args(["--add-links"])
        assert args.add_links is True

    def test_hide_errors_flag(self):
        """--hide-errors sets hide_errors to True."""
        cli = _make_cli()
        args = cli.parser.parse_args(["--hide-errors"])
        assert args.hide_errors is True

    def test_self_flag(self):
        """--self sets self to True."""
        cli = _make_cli()
        args = cli.parser.parse_args(["--self"])
        assert args.self is True

    def test_quiet_flag(self):
        """--quiet sets quiet to True."""
        cli = _make_cli()
        args = cli.parser.parse_args(["--quiet"])
        assert args.quiet is True

    def test_filter_flag(self):
        """--filter stores the supplied filter value."""
        cli = _make_cli()
        args = cli.parser.parse_args(["--filter", "2024"])
        assert args.filter == "2024"

    def test_debug_queries_flag(self):
        """--debug-queries sets debug_queries to True."""
        cli = _make_cli()
        args = cli.parser.parse_args(["--debug-queries"])
        assert args.debug_queries is True


# ---------------------------------------------------------------------------
# _test_str_to_int tests
# ---------------------------------------------------------------------------

class TestStrToInt:
    """Tests for Cli._test_str_to_int helper."""

    def test_integer_string_returns_true(self):
        assert _make_cli()._test_str_to_int("123") is True

    def test_non_integer_string_returns_false(self):
        assert _make_cli()._test_str_to_int("abc") is False

    def test_float_string_returns_false(self):
        assert _make_cli()._test_str_to_int("12.5") is False

    def test_empty_string_returns_false(self):
        assert _make_cli()._test_str_to_int("") is False

    def test_negative_integer_string_returns_true(self):
        assert _make_cli()._test_str_to_int("-5") is True

    def test_zero_string_returns_true(self):
        assert _make_cli()._test_str_to_int("0") is True

    def test_whitespace_string_returns_false(self):
        assert _make_cli()._test_str_to_int("  ") is False


# ---------------------------------------------------------------------------
# _select_slice_from_df tests
# ---------------------------------------------------------------------------

def _make_ref(**kwargs):
    """Build a LookerReference with sensible defaults."""
    return LookerReference(id="1", **kwargs)


class TestSelectSliceFromDf:
    """Tests for Cli._select_slice_from_df."""

    def test_no_label_no_column_returns_dataframe(self):
        """Returns the full DataFrame when neither label nor column is set."""
        cli = _make_cli()
        df = pd.DataFrame({"col1": [1, 2], "col2": [3, 4]})
        result = cli._select_slice_from_df(df, _make_ref())
        assert isinstance(result, pd.DataFrame)

    def test_label_returns_value_from_first_row(self):
        """label selects the named column from the default row 0."""
        cli = _make_cli()
        df = pd.DataFrame({"col1": [10, 20], "col2": [30, 40]})
        result = cli._select_slice_from_df(df, _make_ref(label="col1"))
        assert result == 10

    def test_column_index_returns_value_from_first_row(self):
        """column=1 selects the second column from row 0."""
        cli = _make_cli()
        df = pd.DataFrame({"col1": [10, 20], "col2": [30, 40]})
        result = cli._select_slice_from_df(df, _make_ref(column=1))
        assert result == 30

    def test_row_shifts_selection(self):
        """row=1 makes the slice operate on row index 1."""
        cli = _make_cli()
        df = pd.DataFrame({"col1": [10, 20], "col2": [30, 40]})
        result = cli._select_slice_from_df(df, _make_ref(row=1, label="col1"))
        assert result == 20

    def test_label_takes_priority_over_column(self):
        """When both label and column are supplied, label wins."""
        cli = _make_cli()
        df = pd.DataFrame({"col1": [10, 20], "col2": [30, 40]})
        result = cli._select_slice_from_df(df, _make_ref(label="col1", column=1))
        # label=col1 → 10, not col2 (30)
        assert result == 10

    def test_row_zero_is_default(self):
        """When row is not set the default is index 0."""
        cli = _make_cli()
        df = pd.DataFrame({"val": [99, 1]})
        result = cli._select_slice_from_df(df, _make_ref(label="val"))
        assert result == 99

    def test_column_zero_returns_first_column(self):
        """column=0 returns the first column value."""
        cli = _make_cli()
        df = pd.DataFrame({"a": [7], "b": [8]})
        result = cli._select_slice_from_df(df, _make_ref(column=0))
        assert result == 7


# ---------------------------------------------------------------------------
# _fill_table tests
# ---------------------------------------------------------------------------

def _make_table(rows, cols):
    """Create a python-pptx table with the given dimensions."""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    return slide.shapes.add_table(
        rows, cols, Inches(1), Inches(1), Inches(6), Inches(3)
    ).table


class TestFillTable:
    """Tests for Cli._fill_table."""

    def test_header_row_is_filled(self):
        """Column names appear in the first (header) row."""
        cli = _make_cli()
        table = _make_table(3, 2)
        df = pd.DataFrame({"Name": ["Alice", "Bob"], "Score": [90, 85]})
        cli._fill_table(table, df, headers=True)
        assert table.cell(0, 0).text == "Name"
        assert table.cell(0, 1).text == "Score"

    def test_data_rows_are_filled(self):
        """Data values are written starting at row 1."""
        cli = _make_cli()
        table = _make_table(3, 2)
        df = pd.DataFrame({"Name": ["Alice", "Bob"], "Score": [90, 85]})
        cli._fill_table(table, df, headers=True)
        assert table.cell(1, 0).text == "Alice"
        assert table.cell(1, 1).text == "90"
        assert table.cell(2, 0).text == "Bob"
        assert table.cell(2, 1).text == "85"

    def test_headers_false_leaves_header_row_empty(self):
        """With headers=False the first row is not overwritten."""
        cli = _make_cli()
        table = _make_table(3, 2)
        df = pd.DataFrame({"Name": ["Alice", "Bob"], "Score": [90, 85]})
        cli._fill_table(table, df, headers=False)
        assert table.cell(0, 0).text == ""
        assert table.cell(0, 1).text == ""

    def test_unused_rows_are_cleared(self):
        """Rows beyond the data range are cleared."""
        cli = _make_cli()
        table = _make_table(5, 2)
        # pre-populate a cell that should be cleared afterwards
        table.cell(4, 0).text = "stale"
        df = pd.DataFrame({"Name": ["Alice"], "Score": [90]})
        cli._fill_table(table, df, headers=True)
        assert table.cell(4, 0).text == ""

    def test_extra_df_rows_are_truncated(self):
        """When the DataFrame has more rows than the table, only the first rows fit."""
        cli = _make_cli()
        table = _make_table(2, 2)  # header + 1 data row
        df = pd.DataFrame({"Name": ["Alice", "Bob", "Charlie"], "Score": [90, 85, 78]})
        cli._fill_table(table, df, headers=True)
        assert table.cell(0, 0).text == "Name"
        assert table.cell(1, 0).text == "Alice"

    def test_unused_columns_are_cleared(self):
        """Table columns beyond the DataFrame width are cleared."""
        cli = _make_cli()
        table = _make_table(2, 3)
        table.cell(0, 2).text = "extra"
        df = pd.DataFrame({"Name": ["Alice"], "Score": [90]})
        cli._fill_table(table, df, headers=True)
        assert table.cell(0, 2).text == ""

    def test_values_are_cast_to_string(self):
        """Numeric values are stored as strings in the table."""
        cli = _make_cli()
        table = _make_table(2, 1)
        df = pd.DataFrame({"count": [42]})
        cli._fill_table(table, df, headers=True)
        # Row 0 is the header row; data is always written starting at row 1
        assert table.cell(1, 0).text == "42"


# ---------------------------------------------------------------------------
# Additional _make_df edge cases
# ---------------------------------------------------------------------------

class TestMakeDfEdgeCases:
    """Extra edge-case tests for Cli._make_df."""

    def test_empty_rows_returns_empty_dataframe(self):
        """No data rows → empty DataFrame with no rows."""
        cli = _make_cli()
        result = _make_result(
            dimensions=["view.date"],
            measures=["view.revenue"],
            table_calculations=[],
            rows=[],
        )
        df = cli._make_df(result)
        assert df.empty

    def test_leftover_columns_appended_last(self):
        """Columns absent from metadata appear at the very end."""
        cli = _make_cli()
        result = _make_result(
            dimensions=["view.date"],
            measures=["view.revenue"],
            table_calculations=[],
            rows=[
                {
                    "view.date.value": "2024-01-01",
                    "view.revenue.value": 100,
                    "some_unknown_col": "extra",
                }
            ],
        )
        df = cli._make_df(result)
        cols = list(df.columns)
        assert cols.index("some_unknown_col") > cols.index("revenue")

    def test_single_dimension_column(self):
        """A single dimension produces a one-column DataFrame."""
        cli = _make_cli()
        result = _make_result(
            dimensions=["view.date"],
            measures=[],
            table_calculations=[],
            rows=[{"view.date.value": "2024-01-01"}],
        )
        df = cli._make_df(result)
        assert list(df.columns) == ["date"]

    def test_field_group_variant_used_as_column_name(self):
        """Columns are renamed using field_group_variant, lowercased & spaces → underscores."""
        cli = _make_cli()
        result = _make_result(
            dimensions=[{"name": "view.my_dim", "field_group_variant": "My Dim Label"}],
            measures=[],
            table_calculations=[],
            rows=[{"view.my_dim.value": "test"}],
        )
        df = cli._make_df(result)
        assert "my_dim_label" in df.columns

    def test_multiple_measures_preserve_metadata_order(self):
        """Multiple measures appear in their metadata-declared order."""
        cli = _make_cli()
        result = _make_result(
            dimensions=["view.date"],
            measures=["view.revenue", "view.cost", "view.profit"],
            table_calculations=[],
            rows=[
                {
                    "view.date.value": "2024-01-01",
                    "view.revenue.value": 100,
                    "view.cost.value": 50,
                    "view.profit.value": 50,
                }
            ],
        )
        df = cli._make_df(result)
        cols = list(df.columns)
        assert cols.index("revenue") < cols.index("cost") < cols.index("profit")

    def test_multiple_table_calculations_preserve_metadata_order(self):
        """Table calculations appear in their metadata-declared order."""
        cli = _make_cli()
        result = _make_result(
            dimensions=["view.date"],
            measures=["view.revenue"],
            table_calculations=["calc_a", "calc_b", "calc_c"],
            rows=[
                {
                    "view.date.value": "2024-01-01",
                    "view.revenue.value": 100,
                    "calc_a.value": 1.0,
                    "calc_b.value": 2.0,
                    "calc_c.value": 3.0,
                }
            ],
        )
        df = cli._make_df(result)
        cols = list(df.columns)
        assert cols.index("calc_a") < cols.index("calc_b") < cols.index("calc_c")

    def test_pivots_with_multiple_measures_grouped_by_pivot_then_measure(self):
        """With two measures and two pivot values, columns are grouped by pivot value first."""
        cli = _make_cli()
        result = _make_result(
            dimensions=["view.region"],
            measures=["view.revenue", "view.cost"],
            table_calculations=[],
            rows=[
                {
                    "view.region.value": "North",
                    "view.revenue|FIELD|2024-01.value": 10,
                    "view.cost|FIELD|2024-01.value": 5,
                    "view.revenue|FIELD|2024-02.value": 20,
                    "view.cost|FIELD|2024-02.value": 10,
                }
            ],
            custom_pivots=["view.month"],
            custom_sorts=["view.month asc 0"],
        )
        df = cli._make_df(result)
        cols = list(df.columns)
        revenue_jan = "view.revenue|FIELD|2024-01.value"
        cost_jan = "view.cost|FIELD|2024-01.value"
        revenue_feb = "view.revenue|FIELD|2024-02.value"
        cost_feb = "view.cost|FIELD|2024-02.value"
        assert (
            cols.index(revenue_jan)
            < cols.index(cost_jan)
            < cols.index(revenue_feb)
            < cols.index(cost_feb)
        )

    def test_no_sort_info_pivot_uses_appearance_order(self):
        """When there are no custom_sorts, pivot columns follow Looker's data order."""
        cli = _make_cli()
        result = _make_result(
            dimensions=["view.region"],
            measures=["view.revenue"],
            table_calculations=[],
            rows=[
                {
                    "view.region.value": "North",
                    "view.revenue|FIELD|B.value": 2,
                    "view.revenue|FIELD|A.value": 1,
                }
            ],
            custom_pivots=["view.category"],
            custom_sorts=[],
        )
        df = cli._make_df(result)
        pivot_cols = [c for c in df.columns if "revenue" in c]
        # B appears first in the data → must come first regardless of alpha order
        assert pivot_cols == [
            "view.revenue|FIELD|B.value",
            "view.revenue|FIELD|A.value",
        ]


# ---------------------------------------------------------------------------
# LookerReference model validation tests
# ---------------------------------------------------------------------------

class TestLookerReferenceModel:
    """Tests for LookerReference Pydantic model."""

    def test_integer_id_converted_to_string(self):
        """Integer IDs are coerced to strings by the field validator."""
        ref = LookerReference(id=123)
        assert ref.id == "123"
        assert isinstance(ref.id, str)

    def test_string_id_preserved(self):
        """String IDs pass through unchanged."""
        ref = LookerReference(id="456")
        assert ref.id == "456"

    def test_default_id_type_is_look(self):
        ref = LookerReference(id="1")
        assert ref.id_type == "look"

    def test_default_result_format_is_json_bi(self):
        ref = LookerReference(id="1")
        assert ref.result_format == "json_bi"

    def test_default_apply_formatting_is_false(self):
        ref = LookerReference(id="1")
        assert ref.apply_formatting is False

    def test_default_apply_vis_is_true(self):
        ref = LookerReference(id="1")
        assert ref.apply_vis is True

    def test_default_headers_is_true(self):
        ref = LookerReference(id="1")
        assert ref.headers is True

    def test_default_server_table_calcs_is_true(self):
        ref = LookerReference(id="1")
        assert ref.server_table_calcs is True

    def test_default_retries_is_zero(self):
        ref = LookerReference(id="1")
        assert ref.retries == 0

    def test_default_meta_is_false(self):
        ref = LookerReference(id="1")
        assert ref.meta is False

    def test_default_meta_iterate_is_false(self):
        ref = LookerReference(id="1")
        assert ref.meta_iterate is False

    def test_default_show_latest_chart_label_is_false(self):
        ref = LookerReference(id="1")
        assert ref.show_latest_chart_label is False

    def test_optional_fields_default_to_none(self):
        ref = LookerReference(id="1")
        assert ref.label is None
        assert ref.column is None
        assert ref.row is None
        assert ref.filter is None
        assert ref.filter_overwrites is None
        assert ref.meta_name is None
        assert ref.image_width is None
        assert ref.image_height is None


# ---------------------------------------------------------------------------
# LookerShape model validation tests
# ---------------------------------------------------------------------------

class TestLookerShapeModel:
    """Tests for LookerShape Pydantic model validator."""

    def _base_data(self, shape_type, **overrides):
        data = {
            "shape_id": "0,1",
            "shape_type": shape_type,
            "slide_number": 0,
            "shape_width": 200,
            "shape_height": 100,
            "integration": {"id": "1"},
        }
        data.update(overrides)
        return data

    def test_picture_shape_propagates_image_dimensions(self):
        """PICTURE shapes copy shape_width/height into the integration."""
        shape = LookerShape.model_validate(self._base_data("PICTURE"))
        assert shape.integration.image_width == 200
        assert shape.integration.image_height == 100

    def test_table_shape_sets_apply_formatting_true_by_default(self):
        """TABLE shapes default apply_formatting to True."""
        shape = LookerShape.model_validate(self._base_data("TABLE"))
        assert shape.integration.apply_formatting is True

    def test_table_shape_respects_explicit_apply_formatting_false(self):
        """If apply_formatting is explicitly False on a TABLE, it stays False."""
        data = self._base_data("TABLE")
        data["integration"]["apply_formatting"] = False
        shape = LookerShape.model_validate(data)
        assert shape.integration.apply_formatting is False

    def test_non_picture_shape_has_no_image_dimensions(self):
        """Non-PICTURE shapes do not populate image_width / image_height."""
        shape = LookerShape.model_validate(self._base_data("TEXT_BOX"))
        assert shape.integration.image_width is None
        assert shape.integration.image_height is None

    def test_original_integration_is_preserved(self):
        """original_integration holds a snapshot of the integration before mutation."""
        shape = LookerShape.model_validate(self._base_data("TEXT_BOX"))
        assert shape.original_integration is not None
        assert shape.original_integration.id == "1"

    def test_picture_original_integration_shares_mutations(self):
        """original_integration for PICTURE is assigned by reference before mutation,
        so it also reflects the injected image dimensions (same dict object)."""
        shape = LookerShape.model_validate(self._base_data("PICTURE"))
        # Both integration and original_integration receive the injected dimensions
        # because the validator assigns original_integration = integration (same ref)
        assert shape.original_integration.image_width == 200
        assert shape.original_integration.image_height == 100

    def test_chart_shape_apply_formatting_not_overridden(self):
        """CHART shapes don't have apply_formatting forced to True."""
        shape = LookerShape.model_validate(self._base_data("CHART"))
        assert shape.integration.apply_formatting is False  # default from LookerReference

    def test_shape_id_stored(self):
        shape = LookerShape.model_validate(self._base_data("TEXT_BOX"))
        assert shape.shape_id == "0,1"

    def test_slide_number_stored(self):
        shape = LookerShape.model_validate(self._base_data("TEXT_BOX"))
        assert shape.slide_number == 0


class TestLookerReferenceConfigurationPatterns:
    """Tests that validate the YAML configuration patterns documented in
    docs/getting_started.rst.  Each test corresponds to one documented pattern
    and confirms that LookerReference accepts the configuration without error.
    """

    def test_pattern_simple_id_only(self):
        """Pattern 1 – minimum viable config: only ``id`` is required."""
        ref = LookerReference(id=42)
        assert ref.id == "42"
        assert ref.id_type == "look"

    def test_pattern_row_and_column_selection(self):
        """Pattern 2 – single value extraction by row and column index."""
        ref = LookerReference(id=42, row=0, column=1)
        assert ref.row == 0
        assert ref.column == 1

    def test_pattern_label_selection(self):
        """Pattern 3 – single value extraction by column label."""
        ref = LookerReference(id=42, row=0, label="Total Revenue")
        assert ref.label == "Total Revenue"
        assert ref.row == 0

    def test_pattern_image_result_format(self):
        """Pattern 4 – fetch a Looker chart as a PNG image."""
        ref = LookerReference(id=42, result_format="png")
        assert ref.result_format == "png"

    def test_pattern_image_explicit_dimensions(self):
        """Pattern 4 variant – explicit pixel dimensions for image rendering."""
        ref = LookerReference(id=42, result_format="png", image_width=1200, image_height=675)
        assert ref.image_width == 1200
        assert ref.image_height == 675

    def test_pattern_apply_formatting(self):
        """Pattern 8 – ask Looker to return pre-formatted value strings."""
        ref = LookerReference(id=42, apply_formatting=True)
        assert ref.apply_formatting is True

    def test_pattern_filter_field(self):
        """Pattern 7 – dynamic CLI filter applied to a Looker dimension."""
        ref = LookerReference(id=42, filter="orders.region")
        assert ref.filter == "orders.region"

    def test_pattern_filter_overwrites(self):
        """Pattern 7 variant – static filter overrides baked into the YAML."""
        ref = LookerReference(
            id=42,
            filter_overwrites={"orders.status": "complete", "orders.region": "EMEA"},
        )
        assert ref.filter_overwrites == {"orders.status": "complete", "orders.region": "EMEA"}

    def test_pattern_retries(self):
        """Pattern 9 – retry on transient Looker API failures."""
        ref = LookerReference(id=42, retries=3)
        assert ref.retries == 3

    def test_pattern_id_accepts_integer(self):
        """id field accepts an integer and converts it to a string."""
        ref = LookerReference(id=99)
        assert ref.id == "99"
        assert isinstance(ref.id, str)

    def test_pattern_id_accepts_string(self):
        """id field accepts a string directly."""
        ref = LookerReference(id="99")
        assert ref.id == "99"

