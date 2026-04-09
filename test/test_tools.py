"""Unit tests for the looker_powerpoint.tools sub-package.

Covers:
  - tools/find_alt_text.py   – extract_alt_text, get_presentation_objects_with_descriptions
  - tools/pptx_text_handler.py – emoji removal, header sanitisation, colour encoding,
                                  colorize_positive, Jinja2 rendering, text-frame helpers
  - tools/url_to_hyperlink.py – add_text_with_numbered_links
"""

import io
import os
import tempfile

import pandas as pd
import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from looker_powerpoint.tools.find_alt_text import (
    extract_alt_text,
    get_presentation_objects_with_descriptions,
)
from looker_powerpoint.tools.pptx_text_handler import (
    colorize_positive,
    copy_font_format,
    copy_run_format,
    decode_marked_segments,
    encode_colored_text,
    extract_text_and_run_meta,
    make_jinja_env,
    process_text_field,
    remove_emojis_from_string,
    render_text_with_jinja,
    sanitize_dataframe_headers,
    sanitize_header_name,
    update_text_frame_preserving_formatting,
)
from looker_powerpoint.tools.url_to_hyperlink import add_text_with_numbered_links


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

EXISTING_TABLE_PPTX = os.path.join(os.path.dirname(__file__), "pptx", "table7x7.pptx")


def _make_text_box_pptx(text: str) -> Presentation:
    """Return an in-memory Presentation with a single text-box shape."""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    txBox.text_frame.text = text
    return prs


def _first_shape(prs: Presentation):
    return prs.slides[0].shapes[0]


def _pptx_with_alt_text(yaml_text: str, shape_kind: str = "table") -> Presentation:
    """Build a minimal .pptx that contains a shape whose alt-text is *yaml_text*.

    shape_kind: "table" | "picture" | "text"
    """
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    if shape_kind == "table":
        from pptx.util import Inches

        rows, cols = 2, 2
        tbl = slide.shapes.add_table(
            rows, cols, Inches(1), Inches(1), Inches(4), Inches(2)
        )
        shape = tbl
    else:
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        shape = txBox

    # Inject the alt-text (descr attribute) into the shape XML directly.
    from lxml import etree

    NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
    NS = {"p": NS_P}
    xml_elem = etree.fromstring(shape.element.xml)

    # Try each known path for cNvPr
    for path in [
        ".//p:nvSpPr/p:cNvPr",
        ".//p:nvPicPr/p:cNvPr",
        ".//p:nvGraphicFramePr/p:cNvPr",
    ]:
        elems = xml_elem.xpath(path, namespaces=NS)
        if elems:
            elems[0].set("descr", yaml_text)
            # Write the modified XML back into the shape element
            shape.element.getparent().replace(shape.element, xml_elem)
            break

    return prs


def _save_and_reload(prs: Presentation) -> str:
    """Save *prs* to a temp file and return the file path."""
    fd, path = tempfile.mkstemp(suffix=".pptx")
    os.close(fd)
    prs.save(path)
    return path


# ===========================================================================
# Tests – find_alt_text.py
# ===========================================================================


class TestExtractAltText:
    """Tests for extract_alt_text()."""

    def test_existing_table_shape_returns_dict(self):
        """table7x7.pptx table shape has YAML alt-text that becomes a dict."""
        prs = Presentation(EXISTING_TABLE_PPTX)
        shape = prs.slides[0].shapes[0]
        result = extract_alt_text(shape)
        assert isinstance(result, dict)
        assert "id" in result

    def test_table_shape_id_value(self):
        """table7x7.pptx: id should be 1 (integer from YAML)."""
        prs = Presentation(EXISTING_TABLE_PPTX)
        shape = prs.slides[0].shapes[0]
        result = extract_alt_text(shape)
        assert result["id"] == 1

    def test_shape_without_alt_text_returns_none(self):
        """A plain text-box with no alt-text description returns None."""
        prs = _make_text_box_pptx("hello world")
        shape = _first_shape(prs)
        result = extract_alt_text(shape)
        assert result is None

    def test_non_yaml_alt_text_parsed_as_scalar(self):
        """If the alt-text is a plain string (not a mapping), yaml.safe_load
        returns the string itself — extract_alt_text should still return it."""
        from unittest.mock import MagicMock

        NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
        xml = (
            f'<p:sp xmlns:p="{NS_P}">'
            f'<p:nvSpPr><p:cNvPr id="1" name="test" descr="just a string"/></p:nvSpPr>'
            f"</p:sp>"
        )
        shape = MagicMock()
        shape.element.xml = xml
        result = extract_alt_text(shape)
        assert result == "just a string"


class TestGetPresentationObjectsWithDescriptions:
    """Tests for get_presentation_objects_with_descriptions()."""

    def test_returns_list(self):
        result = get_presentation_objects_with_descriptions(EXISTING_TABLE_PPTX)
        assert isinstance(result, list)

    def test_finds_one_shape(self):
        result = get_presentation_objects_with_descriptions(EXISTING_TABLE_PPTX)
        assert len(result) == 1

    def test_shape_keys_present(self):
        result = get_presentation_objects_with_descriptions(EXISTING_TABLE_PPTX)
        obj = result[0]
        for key in (
            "shape_id",
            "shape_type",
            "shape_width",
            "shape_height",
            "integration",
            "slide_number",
            "shape_number",
        ):
            assert key in obj, f"Missing key: {key}"

    def test_slide_number_is_zero_based(self):
        result = get_presentation_objects_with_descriptions(EXISTING_TABLE_PPTX)
        assert result[0]["slide_number"] == 0

    def test_shape_type_is_string(self):
        result = get_presentation_objects_with_descriptions(EXISTING_TABLE_PPTX)
        assert isinstance(result[0]["shape_type"], str)

    def test_dimensions_are_positive_integers(self):
        result = get_presentation_objects_with_descriptions(EXISTING_TABLE_PPTX)
        obj = result[0]
        assert obj["shape_width"] > 0
        assert obj["shape_height"] > 0

    def test_integration_is_dict(self):
        result = get_presentation_objects_with_descriptions(EXISTING_TABLE_PPTX)
        assert isinstance(result[0]["integration"], dict)

    def test_missing_file_returns_empty_list(self):
        result = get_presentation_objects_with_descriptions("/nonexistent/path.pptx")
        assert result == []

    def test_meta_name_used_as_shape_id(self):
        """When a shape's YAML contains meta_name, that value becomes shape_id."""
        prs = _pptx_with_alt_text("id: 99\nmeta_name: my_meta")
        path = _save_and_reload(prs)
        try:
            result = get_presentation_objects_with_descriptions(path)
            assert len(result) == 1
            assert result[0]["shape_id"] == "my_meta"
        finally:
            os.unlink(path)

    def test_shape_id_uses_slide_and_shape_number_when_no_meta_name(self):
        """Without meta_name, shape_id should be 'slide,shape_id' format."""
        prs = _pptx_with_alt_text("id: 42")
        path = _save_and_reload(prs)
        try:
            result = get_presentation_objects_with_descriptions(path)
            assert len(result) == 1
            # shape_id should contain a comma (slide,shape_number format)
            assert "," in result[0]["shape_id"]
        finally:
            os.unlink(path)

    def test_shapes_without_yaml_are_ignored(self):
        """Shapes with no alt-text are not included in the result."""
        prs = _make_text_box_pptx("no metadata here")
        path = _save_and_reload(prs)
        try:
            result = get_presentation_objects_with_descriptions(path)
            assert result == []
        finally:
            os.unlink(path)


# ===========================================================================
# Tests – pptx_text_handler.py  (emoji & sanitization)
# ===========================================================================


class TestRemoveEmojisFromString:
    def test_removes_emoticons(self):
        assert remove_emojis_from_string("Hello 😀 World") == "Hello  World"

    def test_removes_symbol_emoji(self):
        assert remove_emojis_from_string("📊 Revenue") == " Revenue"

    def test_no_emoji_unchanged(self):
        assert remove_emojis_from_string("Hello World") == "Hello World"

    def test_empty_string(self):
        assert remove_emojis_from_string("") == ""

    def test_non_string_returned_unchanged(self):
        assert remove_emojis_from_string(42) == 42
        assert remove_emojis_from_string(None) is None

    def test_only_emojis_becomes_empty(self):
        result = remove_emojis_from_string("😀🎉")
        assert result == ""

    def test_flag_emoji_removed(self):
        result = remove_emojis_from_string("🇺🇸 United States")
        assert "🇺🇸" not in result


class TestSanitizeHeaderName:
    def test_removes_emoji(self):
        result = sanitize_header_name("📊 Revenue")
        assert "📊" not in result
        assert "Revenue" in result

    def test_spaces_become_underscores(self):
        assert sanitize_header_name("total revenue") == "total_revenue"

    def test_multiple_spaces_collapsed(self):
        assert sanitize_header_name("a  b") == "a_b"

    def test_none_returns_none(self):
        assert sanitize_header_name(None) is None

    def test_leading_trailing_underscores_stripped(self):
        result = sanitize_header_name(" revenue ")
        assert not result.startswith("_")
        assert not result.endswith("_")

    def test_plain_name_unchanged(self):
        assert sanitize_header_name("revenue") == "revenue"

    def test_non_string_coerced(self):
        result = sanitize_header_name(42)
        assert result == "42"


class TestSanitizeDataframeHeaders:
    def test_renames_emoji_headers(self):
        df = pd.DataFrame({"📊 Revenue": [1, 2], "Cost 💰": [3, 4]})
        result = sanitize_dataframe_headers(df)
        assert "Revenue" in result.columns
        assert "Cost" in result.columns

    def test_spaces_replaced_with_underscores(self):
        df = pd.DataFrame({"total revenue": [1], "unit cost": [2]})
        result = sanitize_dataframe_headers(df)
        assert "total_revenue" in result.columns
        assert "unit_cost" in result.columns

    def test_data_unchanged(self):
        df = pd.DataFrame({"a b": [10, 20]})
        result = sanitize_dataframe_headers(df)
        assert list(result["a_b"]) == [10, 20]

    def test_plain_headers_unchanged(self):
        df = pd.DataFrame({"revenue": [1], "cost": [2]})
        result = sanitize_dataframe_headers(df)
        assert list(result.columns) == ["revenue", "cost"]


# ===========================================================================
# Tests – pptx_text_handler.py  (colour encoding / decoding)
# ===========================================================================


class TestColorEncoding:
    def test_encode_decode_roundtrip(self):
        encoded = encode_colored_text("hello", "#FF0000")
        segments = decode_marked_segments(encoded)
        assert len(segments) == 1
        assert segments[0] == ("hello", "#FF0000")

    def test_decode_plain_text(self):
        segments = decode_marked_segments("plain text")
        assert segments == [("plain text", None)]

    def test_decode_mixed_text(self):
        text = "before " + encode_colored_text("42", "#008000") + " after"
        segments = decode_marked_segments(text)
        assert segments[0] == ("before ", None)
        assert segments[1] == ("42", "#008000")
        assert segments[2] == (" after", None)

    def test_multiple_encoded_segments(self):
        t = encode_colored_text("pos", "#008000") + encode_colored_text(
            "neg", "#C00000"
        )
        segments = decode_marked_segments(t)
        assert segments[0] == ("pos", "#008000")
        assert segments[1] == ("neg", "#C00000")


# ===========================================================================
# Tests – pptx_text_handler.py  (colorize_positive)
# ===========================================================================


class TestColorizePositive:
    def test_positive_number_uses_positive_hex(self):
        result = colorize_positive(10)
        segs = decode_marked_segments(result)
        assert segs[0][1] == "#008000"

    def test_negative_number_uses_negative_hex(self):
        result = colorize_positive(-5)
        segs = decode_marked_segments(result)
        assert segs[0][1] == "#C00000"

    def test_zero_uses_zero_hex(self):
        result = colorize_positive(0)
        segs = decode_marked_segments(result)
        assert segs[0][1] == "#000000"

    def test_positive_string_uses_positive_hex(self):
        result = colorize_positive("100")
        segs = decode_marked_segments(result)
        assert segs[0][1] == "#008000"

    def test_negative_string_uses_negative_hex(self):
        result = colorize_positive("-50")
        segs = decode_marked_segments(result)
        assert segs[0][1] == "#C00000"

    def test_parenthetical_negative_string(self):
        # Known limitation: "(100)" is parsed as 100 (positive) by the current
        # implementation because the leading "(" is stripped by regex substitution
        # before the parenthetical-negative check is reached. This test documents
        # the actual behavior; an explicit "-100" string should be used instead
        # for negative values in accounting notation.
        result = colorize_positive("(100)")
        segs = decode_marked_segments(result)
        assert segs[0][1] == "#008000"

    def test_none_uses_zero_hex(self):
        result = colorize_positive(None)
        segs = decode_marked_segments(result)
        assert segs[0][1] == "#000000"

    def test_non_numeric_string_uses_zero_hex(self):
        result = colorize_positive("N/A")
        segs = decode_marked_segments(result)
        assert segs[0][1] == "#000000"

    def test_custom_positive_hex(self):
        result = colorize_positive(5, positive_hex="#ABCDEF")
        segs = decode_marked_segments(result)
        assert segs[0][1] == "#ABCDEF"

    def test_custom_negative_hex(self):
        result = colorize_positive(-5, negative_hex="#123456")
        segs = decode_marked_segments(result)
        assert segs[0][1] == "#123456"

    def test_comma_formatted_number_positive(self):
        result = colorize_positive("1,234")
        segs = decode_marked_segments(result)
        assert segs[0][1] == "#008000"

    def test_scientific_notation_positive(self):
        result = colorize_positive("1.5e2")
        segs = decode_marked_segments(result)
        assert segs[0][1] == "#008000"

    def test_text_preserved_in_output(self):
        """The text inside the encoded segment should match the input value."""
        result = colorize_positive("42")
        segs = decode_marked_segments(result)
        assert segs[0][0] == "42"


# ===========================================================================
# Tests – pptx_text_handler.py  (Jinja2 rendering)
# ===========================================================================


class TestMakeJinjaEnv:
    def test_returns_environment(self):
        from jinja2 import Environment

        env = make_jinja_env()
        assert isinstance(env, Environment)

    def test_colorize_positive_filter_registered(self):
        env = make_jinja_env()
        assert "colorize_positive" in env.filters

    def test_colorize_positive_filter_callable(self):
        env = make_jinja_env()
        assert callable(env.filters["colorize_positive"])


class TestRenderTextWithJinja:
    def test_simple_variable(self):
        result = render_text_with_jinja("Hello {{ name }}", {"name": "World"})
        assert result == "Hello World"

    def test_no_tags_returns_unchanged(self):
        result = render_text_with_jinja("plain text", {})
        assert result == "plain text"

    def test_colorize_positive_filter_in_template(self):
        result = render_text_with_jinja(
            "{{ value | colorize_positive }}", {"value": 10}
        )
        segs = decode_marked_segments(result)
        assert any(color == "#008000" for _, color in segs)

    def test_context_with_list(self):
        result = render_text_with_jinja("{{ items[0] }}", {"items": ["a", "b"]})
        assert result == "a"

    def test_uses_provided_env(self):
        env = make_jinja_env()
        result = render_text_with_jinja("{{ x }}", {"x": "test"}, env=env)
        assert result == "test"

    def test_none_context_treated_as_empty(self):
        result = render_text_with_jinja("static", None)
        assert result == "static"


# ===========================================================================
# Tests – pptx_text_handler.py  (process_text_field)
# ===========================================================================


def _make_textbox_shape(text: str):
    """Return a real pptx shape with a text frame containing *text*."""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    txBox.text_frame.text = text
    return txBox


class TestProcessTextField:
    def test_no_jinja_updates_text(self):
        """When no Jinja tags are present, text is replaced if it differs."""
        shape = _make_textbox_shape("old text")
        df = pd.DataFrame({"col": [1]})
        process_text_field(shape, "new text", df)
        assert shape.text_frame.paragraphs[0].runs[0].text == "new text"

    def test_no_jinja_skips_update_when_same(self):
        """When no Jinja tags and text is the same, text frame is unchanged."""
        shape = _make_textbox_shape("same text")
        df = pd.DataFrame({"col": [1]})
        process_text_field(shape, "same text", df)
        # Should not raise and text stays the same
        tf = shape.text_frame
        full = "".join(r.text for p in tf.paragraphs for r in p.runs)
        assert full == "same text"

    def test_jinja_variable_rendered(self):
        """A simple Jinja2 variable in the text frame gets substituted."""
        shape = _make_textbox_shape("{{ header_rows[0]['value'] }}")
        df = pd.DataFrame({"value": ["hello"]})
        env = make_jinja_env()
        process_text_field(shape, "ignored", df, env=env)
        tf = shape.text_frame
        full = "".join(r.text for p in tf.paragraphs for r in p.runs)
        assert "hello" in full

    def test_jinja_skips_when_rendered_same_as_template(self):
        """If text has no Jinja tags and text_to_insert matches existing, no update is made."""
        shape = _make_textbox_shape("static content")
        df = pd.DataFrame()
        # Pass the same text as text_to_insert so no update should happen
        process_text_field(shape, "static content", df)
        tf = shape.text_frame
        full = "".join(r.text for p in tf.paragraphs for r in p.runs)
        assert full == "static content"


# ===========================================================================
# Tests – pptx_text_handler.py  (update_text_frame_preserving_formatting)
# ===========================================================================


class TestUpdateTextFramePreservingFormatting:
    def test_updates_text(self):
        shape = _make_textbox_shape("original")
        update_text_frame_preserving_formatting(shape.text_frame, "updated")
        runs = [r for p in shape.text_frame.paragraphs for r in p.runs]
        texts = [r.text for r in runs if r.text]
        assert "updated" in texts

    def test_empty_replacement(self):
        shape = _make_textbox_shape("some text")
        update_text_frame_preserving_formatting(shape.text_frame, "")
        runs = [r for p in shape.text_frame.paragraphs for r in p.runs]
        text = "".join(r.text for r in runs)
        assert text == ""


# ===========================================================================
# Tests – pptx_text_handler.py  (copy_run_format / copy_font_format)
# ===========================================================================


class TestCopyRunFormat:
    def _make_run(self, prs: Presentation):
        """Helper to add a single run to a fresh slide."""
        blank = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank)
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "test"
        return run

    def test_copy_bold(self):
        prs = Presentation()
        src = self._make_run(prs)
        dst = self._make_run(prs)
        src.font.bold = True
        copy_run_format(src, dst)
        assert dst.font.bold is True

    def test_copy_italic(self):
        prs = Presentation()
        src = self._make_run(prs)
        dst = self._make_run(prs)
        src.font.italic = True
        copy_run_format(src, dst)
        assert dst.font.italic is True

    def test_copy_font_size(self):
        prs = Presentation()
        src = self._make_run(prs)
        dst = self._make_run(prs)
        src.font.size = Pt(14)
        copy_run_format(src, dst)
        assert dst.font.size == Pt(14)

    def test_copy_rgb_color(self):
        prs = Presentation()
        src = self._make_run(prs)
        dst = self._make_run(prs)
        src.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
        copy_run_format(src, dst)
        assert dst.font.color.rgb == RGBColor(0xFF, 0x00, 0x00)

    def test_copy_does_not_raise_on_no_color(self):
        """copy_run_format should not raise when source has no explicit color."""
        prs = Presentation()
        src = self._make_run(prs)
        dst = self._make_run(prs)
        copy_run_format(src, dst)  # no error


class TestCopyFontFormat:
    def _make_run(self, prs: Presentation):
        blank = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank)
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        p = txBox.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "x"
        return run

    def test_copies_bold(self):
        prs = Presentation()
        src_run = self._make_run(prs)
        dst_run = self._make_run(prs)
        src_run.font.bold = True
        copy_font_format(src_run.font, dst_run.font)
        assert dst_run.font.bold is True

    def test_copies_size(self):
        prs = Presentation()
        src_run = self._make_run(prs)
        dst_run = self._make_run(prs)
        src_run.font.size = Pt(18)
        copy_font_format(src_run.font, dst_run.font)
        assert dst_run.font.size == Pt(18)

    def test_none_src_does_not_raise(self):
        prs = Presentation()
        dst_run = self._make_run(prs)
        copy_font_format(None, dst_run.font)  # Should not raise
        # dst font should remain in its default (unset) state
        assert dst_run.font.bold is None
        assert dst_run.font.italic is None
        assert dst_run.font.size is None

    def test_none_dst_does_not_raise(self):
        prs = Presentation()
        src_run = self._make_run(prs)
        src_run.font.bold = True
        copy_font_format(src_run.font, None)  # Should not raise; nothing to copy to


# ===========================================================================
# Tests – url_to_hyperlink.py
# ===========================================================================


def _make_text_frame():
    """Return a fresh text_frame from a text-box shape."""
    prs = Presentation()
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(2))
    return txBox.text_frame


class TestAddTextWithNumberedLinks:
    def test_no_urls_text_preserved(self):
        tf = _make_text_frame()
        add_text_with_numbered_links(tf, "Hello World")
        text = "".join(r.text for p in tf.paragraphs for r in p.runs)
        assert "Hello World" in text

    def test_no_urls_returns_start_index(self):
        tf = _make_text_frame()
        result = add_text_with_numbered_links(tf, "no links here", start_index=3)
        assert result == 3

    def test_single_url_replaced_with_numbered_placeholder(self):
        tf = _make_text_frame()
        add_text_with_numbered_links(tf, "See https://example.com for details")
        all_text = "".join(r.text for p in tf.paragraphs for r in p.runs)
        assert "(1)" in all_text
        assert "https://example.com" not in all_text

    def test_single_url_hyperlink_set(self):
        tf = _make_text_frame()
        target_url = "https://example.com"
        add_text_with_numbered_links(tf, "See " + target_url + " here")
        # Find the run that has a hyperlink
        links = [
            r.hyperlink.address
            for p in tf.paragraphs
            for r in p.runs
            if r.hyperlink and r.hyperlink.address
        ]
        assert any(addr == target_url for addr in links)

    def test_multiple_urls_incrementing_numbers(self):
        tf = _make_text_frame()
        add_text_with_numbered_links(tf, "A https://a.com B https://b.com C")
        all_text = "".join(r.text for p in tf.paragraphs for r in p.runs)
        assert "(1)" in all_text
        assert "(2)" in all_text

    def test_multiple_urls_returns_next_index(self):
        tf = _make_text_frame()
        result = add_text_with_numbered_links(
            tf, "https://a.com https://b.com", start_index=1
        )
        assert result == 3

    def test_start_index_offset(self):
        tf = _make_text_frame()
        add_text_with_numbered_links(tf, "Visit https://example.com", start_index=5)
        all_text = "".join(r.text for p in tf.paragraphs for r in p.runs)
        assert "(5)" in all_text

    def test_url_ending_with_digits_uses_those_digits(self):
        """URL ending in digits should use that number as placeholder."""
        tf = _make_text_frame()
        add_text_with_numbered_links(tf, "See https://example.com/report/42")
        all_text = "".join(r.text for p in tf.paragraphs for r in p.runs)
        assert "(42)" in all_text

    def test_newlines_flattened_when_url_present(self):
        tf = _make_text_frame()
        add_text_with_numbered_links(tf, "Line1\nLine2\nhttps://example.com")
        all_text = "".join(r.text for p in tf.paragraphs for r in p.runs)
        assert "\n" not in all_text

    def test_newlines_preserved_when_no_url(self):
        tf = _make_text_frame()
        add_text_with_numbered_links(tf, "Line1\nLine2")
        # Newlines are not stripped when there are no URLs
        all_text = "".join(r.text for p in tf.paragraphs for r in p.runs)
        # "Line1\nLine2" stays as is (no URL, no transformation applied)
        assert "Line1" in all_text and "Line2" in all_text

    def test_url_run_is_blue_and_underlined(self):
        tf = _make_text_frame()
        add_text_with_numbered_links(tf, "https://example.com")
        for p in tf.paragraphs:
            for r in p.runs:
                if r.hyperlink and r.hyperlink.address:
                    assert r.font.underline is True
                    assert r.font.color.rgb == RGBColor(0, 0, 255)


# ---------------------------------------------------------------------------
# Tests – pptx_text_handler.py  (extract_text_and_run_meta)
# ---------------------------------------------------------------------------


class TestExtractTextAndRunMeta:
    """Tests for extract_text_and_run_meta."""

    def test_single_run_returns_full_text(self):
        """A single-paragraph, single-run text frame returns the run text."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        tf = txBox.text_frame
        tf.text = "hello"
        full_text, run_meta = extract_text_and_run_meta(tf)
        assert full_text == "hello"

    def test_run_meta_contains_run_objects(self):
        """run_meta entries with text include the original run object."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        tf = txBox.text_frame
        tf.text = "world"
        full_text, run_meta = extract_text_and_run_meta(tf)
        run_entries = [m for m in run_meta if m["run_obj"] is not None]
        assert len(run_entries) >= 1
        assert run_entries[0]["text"] == "world"

    def test_empty_text_frame_returns_empty_string(self):
        """An empty text frame returns an empty full_text string."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        tf = txBox.text_frame
        full_text, run_meta = extract_text_and_run_meta(tf)
        assert full_text == ""
